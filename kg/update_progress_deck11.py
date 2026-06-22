"""
update_progress_deck11.py — Populate slides 8, 9, and 10 of
BCC_progress_meeting_todo.pptx with real values from the 54-row
batch_results.csv (generated 2026-06-15; 10s decision-cycle /
1-cycle lookahead defaults in effect).

  Slide 8  (Simulation results):  refresh 22×8 table + title + footnote
  Slide 9  (WOBJ sweep):          refresh 22×11 table + title + footnote
  Slide 10 (Preview):             populate 23×8 table + title + footnote

Z2 note: wobj_Z2_total = 0.0 for all strategies in this dataset.
'DRL_DENSITY' was absent from the corridor-coordinator build guard in
intersection_controller.py (~lines 18142 & 18532), preventing
notify_bus_granted() from computing offset corrections.
Fixed 2026-06-15. Cells show '0.0*'; footnotes describe the fix.

Note: wobj_Z4_total in the new dataset equals Net_TotalTT_h_Car+Bus+Truck
(raw vehicle-hours), so slide rows 11 and Z4 are numerically identical —
this is correct; the old dataset computed Z4 differently.
"""

import math
import pandas as pd
from pptx import Presentation

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH  = 'batch_results.csv'

RHO_BUS = 40.0
RHO_CAR = 1.5

# ── Formatting helpers ────────────────────────────────────────────────────────
def f1(x):  return f'{x:,.1f}'
def f2(x):  return f'{x:,.2f}'
def fi(x):  return f'{int(round(x)):,}'

FMT = {
    1: f1, 2: f2, 3: f2, 4: f2, 5: f1, 6: fi,  7: f1,  8: f1,
    10: f1, 11: f1, 12: fi, 13: f2, 14: fi, 15: fi, 16: f1,
    18: f1, 19: f1, 20: f1, 21: f1, 22: f1,
}

Z2_FOOTNOTE = (
    '*Z2 (offset-correction objective) = 0.0 in this dataset: '
    "'DRL_DENSITY' was excluded from the corridor-coordinator build guard "
    'in intersection_controller.py (~lines 18142 & 18532), preventing '
    'notify_bus_granted() from computing offset corrections. '
    'Fixed 2026-06-15 — will populate on the next re-run.'
)

# ── Formula computation ───────────────────────────────────────────────────────
def compute_kpi(d):
    """d: DataFrame subset ordered by desired column order.
    Returns dict  formula-key -> pd.Series of values in that order."""
    r = {}
    r[1]  = d['stats_Net_TotalTT_h_Bus']
    r[2]  = d['stats_Net_TotalTT_h_Bus'] * 60.0 / d['stats_N_DistinctBuses']
    r[3]  = d['stats_Net_TotalTT_h_Car']  * 60.0 / d['stats_N_DistinctCars']
    num4  = (d['stats_Net_TotalTT_h_Car']   * RHO_CAR +
             d['stats_Net_TotalTT_h_Bus']   * RHO_BUS +
             d['stats_Net_TotalTT_h_Truck'] * RHO_CAR) * 60.0
    den4  = (d['stats_N_DistinctCars']   * RHO_CAR +
             d['stats_N_DistinctBuses']  * RHO_BUS +
             d['stats_N_DistinctTrucks'] * RHO_CAR)
    r[4]  = num4 / den4
    r[5]  = d['stats_TotalPassDelay_hrs']
    r[6]  = d['stats_PaxEquivPassages']
    r[7]  = d['stats_AvgPassDelay_s']
    r[8]  = d['stats_SidePassDelay_hrs']
    totveh = (d['stats_N_DistinctCars'] +
              d['stats_N_DistinctBuses'] +
              d['stats_N_DistinctTrucks'])
    r[10] = d['stats_Net_Delay_All'] * totveh / 3600.0
    r[11] = (d['stats_Net_TotalTT_h_Car'] +
             d['stats_Net_TotalTT_h_Bus'] +
             d['stats_Net_TotalTT_h_Truck'])
    r[12] = d['stats_Net_TotalDist_Car']
    r[13] = d['stats_Net_TotalDist_Bus']
    r[14] = totveh
    r[15] = d['stats_N_DistinctBuses']
    r[16] = (d['stats_Net_Flow_Car'] +
             d['stats_Net_Flow_Bus'] +
             d['stats_Net_Flow_Truck'])
    r[18] = d['wobj_Z1_total']           # NaN for NO_TSP
    r[19] = d['wobj_Z2_total']           # 0.0 (bug); NaN for NO_TSP
    r[20] = d['wobj_Z3_total']           # NaN for NO_TSP
    r[21] = d['wobj_Z4_total']           # = Net_TotalTT_h sum in new dataset
    r[22] = d['wobj_objective_total']    # NaN for NO_TSP
    return r

# ── Cell-text helper: preserve run formatting, just change text ───────────────
def set_run_text(cell, text):
    p = cell.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for extra in p.runs[1:]:
            extra.text = ''

# ── Data loading ──────────────────────────────────────────────────────────────
df = pd.read_csv(CSV_PATH)

CORE_ORDER = [
    'NO_TSP', 'DCTSP_MARL', 'DCTSP_ZIG',
    'DCTSP_BARGAIN_SPM', 'DCTSP_MP_ECTM', 'DCTSP_BXT',
]
d_core   = (df[df['run_experiment'].isin(CORE_ORDER)]
              .set_index('run_experiment').loc[CORE_ORDER])
kpi_core = compute_kpi(d_core)

WOBJ_ORDER = [
    'WOBJ_SWEEP_A07_B01', 'WOBJ_SWEEP_A07_B02', 'WOBJ_SWEEP_A07_B03',
    'WOBJ_SWEEP_A08_B01', 'WOBJ_SWEEP_A08_B02', 'WOBJ_SWEEP_A08_B03',
    'WOBJ_SWEEP_A09_B01', 'WOBJ_SWEEP_A09_B02', 'WOBJ_SWEEP_A09_B03',
]
d_wobj   = (df[df['run_experiment'].isin(WOBJ_ORDER)]
              .set_index('run_experiment').loc[WOBJ_ORDER])
kpi_wobj = compute_kpi(d_wobj)

prs = Presentation(PPTX_PATH)

# ── Row-to-formula-key maps ───────────────────────────────────────────────────
# Slides 8 & 9 (21 data rows, no Z3):
#   table rows 1-8  → keys 1-8
#   table row  9    → dash, skip
#   table rows 10-16→ keys 10-16
#   table row  17   → dash, skip
#   table row  18   → key 18 (Z1)
#   table row  19   → key 19 (Z2, gets '0.0*')
#   table row  20   → key 21 (Z4)
#   table row  21   → key 22 (objective)
ROW_MAP_89 = {**{i: i for i in range(1, 9)},
              **{i: i for i in range(10, 17)},
              18: 18, 19: 19, 20: 21, 21: 22}

# Slide 10 / Preview (22 data rows, includes Z3):
#   rows 1-8, 10-16 → same as above
#   row 18 → key 18 (Z1)
#   row 19 → key 19 (Z2, gets '0.0*')
#   row 20 → key 20 (Z3)
#   row 21 → key 21 (Z4)
#   row 22 → key 22 (objective)
ROW_MAP_10 = {**{i: i for i in range(1, 9)},
              **{i: i for i in range(10, 17)},
              18: 18, 19: 19, 20: 20, 21: 21, 22: 22}

# NoPriority (col offset 0 = table col 2) rows to leave unchanged (already '—' or '[TBD]')
NOPRIORITY_SKIP_89 = {18, 19, 20, 21}
NOPRIORITY_SKIP_10 = {18, 19, 20, 21, 22}

# ── Slide 8: Simulation results ───────────────────────────────────────────────
s8   = prs.slides[8]
tbl8 = next(sh for sh in s8.shapes if sh.has_table).table

for trow, fkey in ROW_MAP_89.items():
    for ci, exp_idx in enumerate(range(len(CORE_ORDER))):
        col = 2 + ci
        if ci == 0 and trow in NOPRIORITY_SKIP_89:
            continue
        val = kpi_core[fkey].iloc[ci]
        if math.isnan(val):
            continue
        text = '0.0*' if fkey == 19 else FMT[fkey](val)
        set_run_text(tbl8.cell(trow, col), text)

for sh in s8.shapes:
    if sh.has_text_frame and not sh.has_table:
        runs = sh.text_frame.paragraphs[0].runs
        if sh.top < 200000 and len(runs) >= 3:
            runs[2].text = (
                'Simulation results — Kelvin Grove Rd, AM peak: '
                'NoPriority vs. 5 DCTSP strategies '
                '(6 of 54 configs; 10s/1-cycle defaults — current dataset)'
            )
        elif sh.top > 4_000_000 and runs and '*Z2' not in runs[0].text:
            runs[0].text = runs[0].text.rstrip() + '  ' + Z2_FOOTNOTE

# ── Slide 9: WOBJ sweep ───────────────────────────────────────────────────────
s9   = prs.slides[9]
tbl9 = next(sh for sh in s9.shapes if sh.has_table).table

for trow, fkey in ROW_MAP_89.items():
    for ci in range(len(WOBJ_ORDER)):
        col = 2 + ci
        val = kpi_wobj[fkey].iloc[ci]
        if math.isnan(val):
            continue
        text = '0.0*' if fkey == 19 else FMT[fkey](val)
        set_run_text(tbl9.cell(trow, col), text)

for sh in s9.shapes:
    if sh.has_text_frame and not sh.has_table:
        runs = sh.text_frame.paragraphs[0].runs
        if sh.top < 200000 and len(runs) >= 3:
            runs[2].text = (
                'Sensitivity of objectives — WOBJ_ALPHA × WOBJ_BETA sweep '
                '(GLOBAL_REWARD+KALMAN; 9 of 54 configs; 10s/1-cycle defaults)'
            )
        elif sh.top > 4_000_000 and runs and '*Z2' not in runs[0].text:
            runs[0].text = runs[0].text.rstrip() + '  ' + Z2_FOOTNOTE

# ── Slide 10: Preview ─────────────────────────────────────────────────────────
s10   = prs.slides[10]
tbl10 = next(sh for sh in s10.shapes if sh.has_table).table

for trow, fkey in ROW_MAP_10.items():
    for ci in range(len(CORE_ORDER)):
        col = 2 + ci
        if ci == 0 and trow in NOPRIORITY_SKIP_10:
            continue
        val = kpi_core[fkey].iloc[ci]
        if math.isnan(val):
            continue
        text = '0.0*' if fkey == 19 else FMT[fkey](val)
        set_run_text(tbl10.cell(trow, col), text)

NEW_FOOTNOTE_10 = (
    'Values from the 54-config re-run (2026-06-15; 10s decision-recalculation cycle / '
    '1-cycle lookahead defaults now in effect). '
    "NoPriority Z3 cell: '[TBD]' — bus-lateness is not computed for the no-TSP "
    'baseline in this dataset. '
    "Rows marked '—' (side-street split, throughput-side, and NoPriority's "
    'Z1/Z2/Z4/objective cells) are not applicable / not computable from current outputs. '
    + Z2_FOOTNOTE
)

for sh in s10.shapes:
    if sh.has_text_frame and not sh.has_table:
        runs = sh.text_frame.paragraphs[0].runs
        if sh.top < 200000 and len(runs) >= 3:
            runs[2].text = (
                'Preview — full KPI breakdown incl. new Bus-lateness (Z3) row, '
                'populated for 6 of 54 configs '
                '(10s / 1-cycle defaults now in effect)'
            )
        elif sh.top > 4_000_000 and runs:
            runs[0].text = NEW_FOOTNOTE_10

prs.save(PPTX_PATH)
print('Saved', PPTX_PATH)
