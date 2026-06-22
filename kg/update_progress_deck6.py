"""One-off follow-up script (run after update_progress_deck5.py):

  - Slide "1.4 - Future work": the title (cloned from slide 1.2's 32pt title)
    overflowed off the right edge of the slide, and the table's last row
    overlapped the footer band. Fix: shrink the title to 24pt, shorten the
    longest cell's text slightly, shrink table body/header fonts to 9/10pt,
    and reduce the table's height/position so it clears the footer.

  - New "Preview" slide: the cross-slide XML clone of slide 7's title/table
    rendered badly (the title inherited a stray yellow <a:highlight> from the
    cloned run, and the CellSearch/CellQLearn header columns lost their fill).
    Fix: remove the broken title/table/footnote and rebuild them from scratch
    with add_textbox()/add_table(), matching slide 7/8's column widths, header
    fill (#00467F, white bold text) and NoPriority-column highlight
    (#FFF2CC), with no highlight on the title.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'

HEADER_FILL = RGBColor(0x00, 0x46, 0x7F)
NOPRIORITY_FILL = RGBColor(0xFF, 0xF2, 0xCC)
TITLE_COLOR = RGBColor(0x00, 0x46, 0x7F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(PPTX_PATH)


def _set_cell_lines(cell, lines, size_pt=10, bold=False, color=None):
    tf = cell.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = lines[0]
    p0.runs[0].font.size = Pt(size_pt)
    p0.runs[0].font.bold = bold
    if color is not None:
        p0.runs[0].font.color.rgb = color
    for line in lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def _set_cell(cell, text, size_pt=10, bold=False, color=None):
    _set_cell_lines(cell, [text], size_pt=size_pt, bold=bold, color=color)


# =============================================================================
# Fix 1: Slide 6 (0-indexed) -- "1.4  Future work"
# =============================================================================
s6 = prs.slides[6]
shapes6 = {sh.name: sh for sh in s6.shapes}

title6 = shapes6['TextBox 1']
for p in title6.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(24)

tbl6_shape = shapes6['Table 301']
tbl6 = tbl6_shape.table

# Shrink fonts: header 10pt, body 9pt.
for ri in range(len(tbl6.rows)):
    for ci in range(len(tbl6.columns)):
        for p in tbl6.cell(ri, ci).text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10 if ri == 0 else 9)

# Shorten the longest description (row 1, Description col) slightly.
_set_cell_lines(
    tbl6.cell(1, 1),
    ["Re-centre the existing 90–180s × 1.0–2.0-cycle grid around the new "
     "10s / 1-cycle baseline (Slide 1.3); also revisit the 30s floors in "
     "_reward_future_term for short cycles."],
    size_pt=9, bold=False,
)

# Reduce table height/position so it clears the footer band.
tbl6_shape.top = Inches(1.4)
tbl6_shape.height = Inches(5.0)
row_heights6 = [0.4, 0.92, 0.92, 0.92, 0.92, 0.92]
for ri, h in enumerate(row_heights6):
    tbl6.rows[ri].height = Inches(h)

# =============================================================================
# Fix 2: Slide 10 (0-indexed) -- "Preview" slide -- rebuild title/table/
# footnote from scratch (the cross-slide XML clone rendered incorrectly).
# =============================================================================
s10 = prs.slides[10]
shapes10 = {sh.name: sh for sh in s10.shapes}

for name in ('title', 'table', 'TextBox 14'):
    sh = shapes10[name]
    sh._element.getparent().remove(sh._element)

# -- Title --
title_box = s10.shapes.add_textbox(Inches(0.328), Inches(0.067), Inches(12.133), Inches(1.137))
tf_title = title_box.text_frame
tf_title.word_wrap = True
p_title = tf_title.paragraphs[0]
for text in (
    "2  · ",
    "Project Update and Findings: ",
    "Preview — target table structure for full results "
    "(101 configs, post re-run with 10s / 1-cycle defaults)",
):
    r = p_title.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

# -- Table (23 rows x 8 cols, matching slide 8's "Simulation results" layout) --
N_ROWS, N_COLS = 23, 8
tbl_shape = s10.shapes.add_table(N_ROWS, N_COLS, Inches(1.021), Inches(0.682), Inches(11.6), Inches(6.21))
tbl = tbl_shape.table

tbl.columns[0].width = Inches(2.70)
tbl.columns[1].width = Inches(0.85)
for ci in range(2, 8):
    tbl.columns[ci].width = Inches(1.3417)

row_h = Inches(6.21) / N_ROWS
for ri in range(N_ROWS):
    tbl.rows[ri].height = int(row_h)

headers = ['KPI', 'unit', 'NoPriority', 'CPD-QL', 'WaveGate', 'NashGate', 'CellSearch', 'CellQLearn']
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_FILL
    _set_cell(cell, h, size_pt=9, bold=True, color=WHITE)

TBD = "[TBD]"
TBD6 = [TBD] * 6
DASH6 = ["—"] * 6
DASH_NP = ["—"] + [TBD] * 5  # NoPriority "--" (Z1/Z2/Z4/Total: not applicable)

rows_data = [
    ("Total bus travel time", "[bus-h]", TBD6),
    ("Average bus travel time", "[min/bus]", TBD6),
    ("Total car travel time", "[min/car]", TBD6),
    ("Total passenger travel time", "[min/pax]", TBD6),
    ("Total passenger delay", "[pax-h]", TBD6),
    ("Total serviced passengers", "[pax]", TBD6),
    ("Average passenger delay", "[sec/pax]", TBD6),
    ("Side-street total passenger delay", "[pax-h]", TBD6),
    ("Side-street average passenger delay", "[sec/pax]", DASH6),
    ("Total system delay", "[veh-h]", TBD6),
    ("Total hours of travel", "[veh-h]", TBD6),
    ("Total VKT – car", "[veh-km]", TBD6),
    ("Total VKT – bus", "[veh-km]", TBD6),
    ("Total vehicles in system", "[veh (count)]", TBD6),
    ("Total vehicles in system (Bus)", "[veh (count)]", TBD6),
    ("Throughput – main", "[veh/h]", TBD6),
    ("Throughput – side", "[veh/h]", DASH6),
    ("Weighted pax-delay objective (Z1)", "[a.u.]", DASH_NP),
    ("Offset-correction objective (Z2)", "[a.u.]", DASH_NP),
    ("Bus lateness (Z3, Σσ⁺)", "[bus-s]", TBD6),
    ("Raw travel-time objective (Z4)", "[a.u.]", DASH_NP),
    ("Total weighted objective", "[a.u.]", DASH_NP),
]
assert len(rows_data) == N_ROWS - 1

for ri, (label, unit, vals) in enumerate(rows_data, start=1):
    _set_cell(tbl.cell(ri, 0), label, size_pt=9, bold=False)
    _set_cell(tbl.cell(ri, 1), unit, size_pt=9, bold=False)
    for ci, v in enumerate(vals, start=2):
        _set_cell(tbl.cell(ri, ci), v, size_pt=9, bold=False)
    npcell = tbl.cell(ri, 2)
    npcell.fill.solid()
    npcell.fill.fore_color.rgb = NOPRIORITY_FILL

# -- Footnote --
fn_box = s10.shapes.add_textbox(Inches(1.021), Inches(6.95), Inches(11.6), Inches(0.45))
tf_fn = fn_box.text_frame
tf_fn.word_wrap = True
r = tf_fn.paragraphs[0].add_run()
r.text = (
    "[TBD] = to be populated once the full re-run (101 configurations; 10s "
    "decision-recalculation cycle / 1-cycle lookahead defaults, Slide 1.3) "
    "completes. New row ‘Bus lateness (Z3, Σσ⁺)’ is "
    "populated for ALL strategies including NoPriority — enabling a direct "
    "bus-lateness comparison against the no-TSP baseline (Slide 1.4, Future "
    "work). Rows marked ‘—’ (side-street split, throughput-side, "
    "and NoPriority's Z1/Z2/Z4/objective cells) remain not applicable / not "
    "computable from current outputs, as on the preceding slides."
)
r.font.size = Pt(9)
r.font.italic = True

# =============================================================================
# Save
# =============================================================================
prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
