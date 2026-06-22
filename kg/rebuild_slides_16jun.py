"""
rebuild_slides_16jun.py — Rebuild sensitivity slides in BCC_progress_meeting_16Jun_update.pptx
with NO_TSP side-by-side comparison for every variant.

  Slide 8 (Demand): 22r x 10c — KPI | NT 0.8 | Nash 0.8 | D% | NT 1.0 | Nash 1.0 | D% | NT 1.2 | Nash 1.2 | D%
  Slide 9 (Occ):    22r x 10c — KPI | NT vs LOW | D% | NT vs BASE | D% | NT vs HIGH | D%
                     (each level has own NO_TSP + NashGate + delta)
  Slide 10 (Wobj):  13r x 12c — KPI | NO_TSP | combo1 | D% | combo2 | D% | ... | combo5 | D%
"""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

PPTX_PATH = 'BCC_progress_meeting_16Jun_update.pptx'
HEADER_BG = '2F5496'
HEADER_FG = 'FFFFFF'

def set_cell(tbl, r, c, text, bold=False, sz=Pt(7), color=None):
    cell = tbl.cell(r, c)
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(text); run.font.size = sz; run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(int(color[0:2],16),int(color[2:4],16),int(color[4:6],16))

def hdr_fill(tbl):
    for ci in range(len(tbl.columns)):
        tc = tbl.cell(0, ci)._tc.get_or_add_tcPr()
        sf = etree.SubElement(tc, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr').set('val', HEADER_BG)

def remove_tables(slide):
    for sh in list(slide.shapes):
        if sh.has_table: sh._element.getparent().remove(sh._element)

LABELS_22 = [
    'Bus TT (h)','Avg Bus TT (min)','Avg Car TT (min)','Wtd Avg TT (min)',
    'Total Pass Delay (h)','Pax-Equiv Passages','Avg Pass Delay (s)',
    'Side Pass Delay (h)','---','Vehicle Delay (h)',
    'Total Travel Time (h)','Car Dist (km)','Bus Dist (km)',
    'Total Vehicles','Distinct Buses','Flow (veh/h)',
    '---','Z1: pax delay (pax s)','Z2: offset (s)','Z4: corridor TT (veh h)','Obj = aZ1+bZ2+gZ3',
]
LABELS_WOBJ = [
    'Total Pass Delay (h)','Avg Pass Delay (s)','Side Pass Delay (h)',
    'Total Travel Time (h)','Total Vehicles','Distinct Buses','Flow (veh/h)',
    'Z1: pax delay (pax s)','Z2: offset (s)','Z3: lateness s+ (s)','Z4: corridor TT (veh h)','Obj = aZ1+bZ2+gZ3+dZ4',
]

prs = Presentation(PPTX_PATH)
LEFT, TOP, W, H = Inches(0.4), Inches(0.55), Inches(11.2), Inches(4.0)

# Find slides by title
d_idx = o_idx = w_idx = None
for i, s in enumerate(prs.slides):
    txt = ' '.join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
    if 'Demand scaling' in txt: d_idx = i
    if 'Occupancy' in txt: o_idx = i
    if 'Objective weight' in txt: w_idx = i

print(f'Slides: demand={d_idx} occ={o_idx} wobj={w_idx}')

# --- Slide: Demand (22r x 10c) ---
if d_idx is not None:
    s = prs.slides[d_idx]
    remove_tables(s)
    tbl = s.shapes.add_table(22, 10, LEFT, TOP, W, H).table
    cw = [1.6] + [0.95]*9
    for ci, w_ in enumerate(cw): tbl.columns[ci].width = Inches(w_)
    hdr_fill(tbl)
    set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG)
    for gi, lbl in enumerate(['0.8x','1.0x','1.2x']):
        bc = 1 + gi*3
        set_cell(tbl, 0, bc,   f'NO_TSP {lbl}', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, bc+1, f'NashGate {lbl}', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, bc+2, 'D%', bold=True, color=HEADER_FG)
    for ri, lb in enumerate(LABELS_22):
        set_cell(tbl, ri+1, 0, lb, bold=('---' in lb), sz=Pt(7))
    print('  Demand: 22r x 10c')

# --- Slide: Occupancy (22r x 10c) — 3 groups of (NO_TSP + NashGate + D%) ---
if o_idx is not None:
    s = prs.slides[o_idx]
    remove_tables(s)
    tbl = s.shapes.add_table(22, 10, LEFT, TOP, W, H).table
    cw = [1.6] + [0.95]*9
    for ci, w_ in enumerate(cw): tbl.columns[ci].width = Inches(w_)
    hdr_fill(tbl)
    set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG)
    for gi, lbl in enumerate(['LOW\n(20/1.0)','BASE\n(40/1.2)','HIGH\n(60/1.5)']):
        bc = 1 + gi*3
        set_cell(tbl, 0, bc,   f'NO_TSP', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, bc+1, f'NashGate {lbl}', bold=True, color=HEADER_FG, sz=Pt(7))
        set_cell(tbl, 0, bc+2, 'D%', bold=True, color=HEADER_FG)
    for ri, lb in enumerate(LABELS_22):
        set_cell(tbl, ri+1, 0, lb, bold=('---' in lb), sz=Pt(7))
    print('  Occupancy: 22r x 10c')

# --- Slide: Wobj (13r x 12c) — KPI + NO_TSP + 5 combos x (value + D%) ---
if w_idx is not None:
    s = prs.slides[w_idx]
    remove_tables(s)
    tbl = s.shapes.add_table(13, 12, LEFT, TOP, W, H).table
    cw = [1.4] + [0.9]*11
    for ci, w_ in enumerate(cw): tbl.columns[ci].width = Inches(w_)
    hdr_fill(tbl)
    set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG, sz=Pt(7))
    set_cell(tbl, 0, 1, 'NO_TSP', bold=True, color=HEADER_FG, sz=Pt(7))
    wh = ['EQ Z1+Z2','D%','EQ ALL3','D%','Only Z1','D%','Only Z2','D%','Only TT','D%']
    for ci, hdr in enumerate(wh):
        set_cell(tbl, 0, ci+2, hdr, bold=True, color=HEADER_FG, sz=Pt(5.5))
    for ri, lb in enumerate(LABELS_WOBJ):
        set_cell(tbl, ri+1, 0, lb, sz=Pt(6))
    print('  Wobj: 13r x 12c')

prs.save(PPTX_PATH)
print(f'Saved {PPTX_PATH}')
