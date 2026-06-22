"""One-off follow-up script (run after update_progress_deck2.py):

  - Slides "2 - ... Simulation results" and "2 - ... Sensitivity of objectives
    (WOBJ_ALPHA x WOBJ_BETA)" (the 21-row KPI tables):
      * Insert a new "Total bus travel time [bus-h]" (TBTT) row, per
        KPIS___113.pdf, and relabel the existing "Total bus travel time
        [min/bus]" row -> "Average bus travel time [min/bus]" (ABTT) -- it
        was already computed with the ABTT formula.
      * Relabel "Total vehicles in system[(Bus)]" units to "[veh (count)]"
        to match the PDF's unit notation.
      * Highlight the NoPriority column (no-TSP baseline) on the Simulation
        results table and add footnotes clarifying NoPriority (no-TSP
        baseline) vs TSP-enabled columns on both tables.
  - Slide "1.2 - Sensitivity sweeps configured this period": append a
    sentence about the two new sweeps (decision-cycle/lookahead-horizon and
    the NashGate-restricted weighted-objective alpha/beta/gamma sweep) and
    update the total configured-experiment count.
  - New slide "1.3 - Decision-cycle & lookahead-horizon configuration":
    discloses CycleTime (decision-recalculation cycle, default 135s) and
    REWARD_FUTURE_HORIZON_CYCLES (lookahead horizon, default 1.5 cycles =
    202.5s), with a 4x4 sensitivity grid for the new CYCLE_HORIZON_SWEEP
    (NashGate-only), plus a note on WOBJ_GAMMA_SWEEP_NASHGATE.

Source data: batch_results.csv (for the new TBTT row).
"""

from copy import deepcopy

import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH = 'batch_results.csv'

NOPRIORITY_FILL = RGBColor(0xFF, 0xF2, 0xCC)  # light amber: no-TSP baseline column

# =============================================================================
# 1. Load data -- TBTT (Total bus travel time, [bus-h]) for the new row
# =============================================================================
df = pd.read_csv(CSV_PATH)

STRATS = [
    ("NoPriority", "NO_TSP"),
    ("CPD-QL",     "DCTSP_MARL"),
    ("WaveGate",   "DCTSP_ZIG"),
    ("NashGate",   "DCTSP_BARGAIN_SPM"),
    ("CellSearch", "DCTSP_MP_ECTM"),
    ("CellQLearn", "DCTSP_BXT"),
]
WOBJ_ALPHAS = [0.7, 0.8, 0.9]
WOBJ_BETAS = [0.1, 0.2, 0.3]


def fmt(v, nd=1):
    if pd.isna(v):
        return "—"
    return f"{v:,.{nd}f}"


tbtt_strats = []
for _, exp in STRATS:
    r = df[df['run_experiment'] == exp].iloc[0]
    tbtt_strats.append(fmt(r["stats_Net_TotalTT_h_Bus"], 1))

tbtt_wobj = []
for a in WOBJ_ALPHAS:
    for b in WOBJ_BETAS:
        name = f"WOBJ_SWEEP_A{a:.1f}_B{b:.1f}".replace(".", "")
        r = df[df['run_experiment'] == name].iloc[0]
        tbtt_wobj.append(fmt(r["stats_Net_TotalTT_h_Bus"], 1))

# =============================================================================
# 2. Open presentation
# =============================================================================
prs = Presentation(PPTX_PATH)


def _set_cell_lines(cell, lines, size_pt=9, bold=False):
    """Set a table cell's text to one paragraph per entry in `lines`,
    applying explicit font size/bold to every run."""
    tf = cell.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = lines[0]
    p0.runs[0].font.size = Pt(size_pt)
    p0.runs[0].font.bold = bold
    for line in lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def _set_cell(cell, text, size_pt=9, bold=False):
    _set_cell_lines(cell, [text], size_pt=size_pt, bold=bold)


def _insert_tbtt_row(tbl_shape, tbtt_values, size_pt, n_rows):
    """Insert a new 'Total bus travel time [bus-h]' (TBTT) row right after
    the header, by cloning the existing row 1 (which holds the ABTT data),
    then relabel the now-duplicated original row 1 -> 'Average bus travel
    time'. Returns the re-fetched table (now n_rows+1 rows)."""
    tbl = tbl_shape.table
    _tbl = tbl._tbl
    tr_lst = _tbl.tr_lst
    new_tr = deepcopy(tr_lst[1])
    tr_lst[1].addprevious(new_tr)

    tbl = tbl_shape.table  # re-fetch after XML mutation
    n_rows_new = len(tbl.rows)
    assert n_rows_new == n_rows + 1

    # Even row heights across the (now taller) table.
    total_h = Inches(6.21)
    for row in tbl.rows:
        row.height = int(total_h / n_rows_new)

    # New row 1: TBTT.
    _set_cell(tbl.cell(1, 0), "Total bus travel time", size_pt=size_pt, bold=False)
    _set_cell(tbl.cell(1, 1), "[bus-h]", size_pt=size_pt, bold=False)
    for ci, val in enumerate(tbtt_values, start=2):
        _set_cell(tbl.cell(1, ci), val, size_pt=size_pt, bold=False)

    # Row 2 (originally row 1): relabel to ABTT -- values/unit unchanged.
    _set_cell(tbl.cell(2, 0), "Average bus travel time", size_pt=size_pt, bold=False)

    return tbl


def _relabel_veh_count_units(tbl):
    """Match KPIS___113.pdf's '[veh (count)]' unit notation for the two
    throughput-population rows."""
    for ri in range(len(tbl.rows)):
        label = tbl.cell(ri, 0).text_frame.text
        if label in ("Total vehicles in system", "Total vehicles in system (Bus)"):
            size_pt = tbl.cell(ri, 1).text_frame.paragraphs[0].runs[0].font.size
            size_pt = 9 if size_pt is None else round(size_pt.pt)
            _set_cell(tbl.cell(ri, 1), "[veh (count)]", size_pt=size_pt, bold=False)


def _add_footnote(slide, text, top_in, left_in=1.021, width_in=11.6, size_pt=9):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.45))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.italic = True
    return box


# =============================================================================
# 3. Slide 6 (0-indexed) -- "Simulation results" KPI table (8 cols)
# =============================================================================
s6 = prs.slides[6]
shapes6 = {sh.name: sh for sh in s6.shapes}
tbl6_shape = shapes6["table"]

tbl6 = _insert_tbtt_row(tbl6_shape, tbtt_strats, size_pt=9, n_rows=21)
_relabel_veh_count_units(tbl6)

# Highlight the NoPriority column (col 2) -- the no-TSP baseline -- across all
# data rows (row 0 is the header and keeps its existing header fill/style).
for ri in range(1, len(tbl6.rows)):
    cell = tbl6.cell(ri, 2)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NOPRIORITY_FILL

_add_footnote(
    s6,
    "NoPriority = no-TSP baseline (fixed-time signals, no bus detection or "
    "priority; highlighted column). CPD-QL / WaveGate / NashGate / CellSearch "
    "/ CellQLearn = 5 coordinated DCTSP (dynamic TSP) strategies (Slide 1.1).",
    top_in=6.95,
)

# =============================================================================
# 4. Slide 7 (0-indexed) -- "Sensitivity of objectives (WOBJ_ALPHA x WOBJ_BETA)"
#    KPI table (11 cols)
# =============================================================================
s7 = prs.slides[7]
shapes7 = {sh.name: sh for sh in s7.shapes}
tbl7_shape = shapes7["table"]

tbl7 = _insert_tbtt_row(tbl7_shape, tbtt_wobj, size_pt=8, n_rows=21)
_relabel_veh_count_units(tbl7)

_add_footnote(
    s7,
    "All 9 configurations enable TSP (GLOBAL_REWARD coordination, KALMAN "
    "algorithm; α/β = WOBJ_ALPHA/WOBJ_BETA weights on Z1/Z2). For the "
    "no-TSP baseline (NoPriority), see the NoPriority column on the previous slide.",
    top_in=6.95,
)

# =============================================================================
# 5. Slide 4 (0-indexed) -- "1.2 Sensitivity sweeps configured this period":
#    extend the caption with the two new sweeps + updated total count.
# =============================================================================
s4 = prs.slides[4]
shapes4 = {sh.name: sh for sh in s4.shapes}
caption = shapes4["TextBox 3"]
cap_run = caption.text_frame.paragraphs[0].runs[0]
cap_run.text = cap_run.text.replace(
    "Total configured experiment set this period: 65 (+ 6 demand-mode runs when enabled).",
    "Two further sweeps target the decision-recalculation cycle and reward "
    "lookahead horizon, and the weighted-objective α/β/γ weights "
    "-- both restricted to NashGate (BARGAIN_SPM) and configured, pending run "
    "(see Slide 1.3 for the cycle/horizon disclosure and sensitivity grid: "
    "9 + 27 = 36 further configurations). Total configured experiment set "
    "this period: 101 (+ 6 demand-mode runs when enabled)."
)

# =============================================================================
# 6. New slide "1.3 - Decision-cycle & lookahead-horizon configuration"
# =============================================================================
new_slide = prs.slides.add_slide(prs.slide_layouts[8])  # "3_Standard slide"

for shape in list(new_slide.shapes):
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
        shape._element.getparent().remove(shape._element)

# -- Title: clone slide 1.2's title textbox and edit the text. --
title4 = shapes4["TextBox 1"]
new_title_el = deepcopy(title4._element)
new_slide.shapes._spTree.append(new_title_el)
new_title = new_slide.shapes[-1]
cNvPr = new_title._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '201')
new_title.text_frame.paragraphs[0].runs[0].text = (
    "1.3  ·  Decision-cycle & lookahead-horizon configuration"
)

# -- Disclosure bullets --
bullets = new_slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.13), Inches(2.3))
btf = bullets.text_frame
btf.word_wrap = True
bullet_items = [
    ("Decision-recalculation cycle (CycleTime / TSP_CYCLE_LENGTH_OVERRIDE_S): ",
     "each junction's TSP controller makes ONE grant/no-action decision per "
     "approaching bus within this window. Default: 135s per junction "
     "(intersection_configs.py); now overridable per-run for sensitivity testing."),
    ("Lookahead horizon (REWARD_FUTURE_HORIZON_CYCLES × cycle length): ",
     "the reward's future-debt term projects how much bus/cross-traffic delay "
     "disturbance is likely to need correcting this many cycles ahead. "
     "Default: 1.5 cycles = 1.5 × 135s = 202.5s ahead."),
    ("Sensitivity (CYCLE_HORIZON_SWEEP, NashGate / BARGAIN_SPM target): ",
     "9 configurations crossing 3 cycle lengths × 3 horizon multipliers "
     "(grid below) — configured, pending run."),
]
for i, (lead, rest) in enumerate(bullet_items):
    p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
    r1 = p.add_run()
    r1.text = "•  " + lead
    r1.font.size = Pt(12)
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(12)
    r2.font.bold = False

# -- Sensitivity grid: horizon_s = cycle_s x horizon_cycles --
CYCLES = [90.0, 135.0, 180.0]
HORIZONS = [1.0, 1.5, 2.0]
grid = new_slide.shapes.add_table(4, 4, Inches(0.6), Inches(3.75), Inches(12.13), Inches(2.0))
gtbl = grid.table
gtbl.columns[0].width = Inches(3.5)
for ci in range(1, 4):
    gtbl.columns[ci].width = Inches(2.877)
for ri in range(4):
    gtbl.rows[ri].height = Inches(0.5)

_set_cell(gtbl.cell(0, 0), "Recalculation cycle ↓  /  Lookahead horizon →", size_pt=11, bold=True)
for ci, h in enumerate(HORIZONS, start=1):
    label = f"{h:.1f} × cycle" + (" (default)" if h == 1.5 else "")
    _set_cell(gtbl.cell(0, ci), label, size_pt=11, bold=True)

for ri, c in enumerate(CYCLES, start=1):
    row_label = f"{c:.0f} s" + (" (default)" if c == 135.0 else "")
    _set_cell(gtbl.cell(ri, 0), row_label, size_pt=11, bold=True)
    for ci, h in enumerate(HORIZONS, start=1):
        horizon_s = c * h
        is_default = (c == 135.0 and h == 1.5)
        label = f"{horizon_s:.1f} s" + (" (default)" if is_default else "")
        _set_cell(gtbl.cell(ri, ci), label, size_pt=11, bold=is_default)
        if is_default:
            cell = gtbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)

# -- Closing note: WOBJ_GAMMA_SWEEP_NASHGATE --
note = new_slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.13), Inches(1.2))
ntf = note.text_frame
ntf.word_wrap = True
r = ntf.paragraphs[0].add_run()
r.text = (
    "•  Weighted-objective α/β/γ sensitivity "
    "(WOBJ_GAMMA_SWEEP_NASHGATE): J = WOBJ_ALPHA·Z1 + WOBJ_BETA·Z2 + "
    "WOBJ_GAMMA·Z3, where Z1 = weighted pax-delay, Z2 = offset-correction "
    "error, Z3 = bus lateness (Σσ⁺). 27 configurations "
    "(WOBJ_ALPHA ∈ {0.7,0.8,0.9} × WOBJ_BETA ∈ {0.1,0.2,0.3} × "
    "WOBJ_GAMMA ∈ {0.0,0.1,0.2}), restricted to NashGate (BARGAIN_SPM, the "
    "best-performing strategy) — configured, pending run."
)
r.font.size = Pt(12)

# Reposition the new slide immediately after slide 1.2 (index 4 -> 5).
sldIdLst = prs.slides._sldIdLst
new_sldId = sldIdLst[-1]
sldIdLst.remove(new_sldId)
sldIdLst.insert(5, new_sldId)

# =============================================================================
# Save
# =============================================================================
prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
print(f"Total slides: {len(prs.slides)}")
