"""One-off follow-up script (run after update_progress_deck4.py):

  - Slide "1.3 - Decision-cycle & lookahead-horizon configuration": simplify.
    Removes the 4x4 CYCLE_HORIZON_SWEEP sensitivity grid and the
    WOBJ_GAMMA_SWEEP_NASHGATE closing note, and rewrites the 3 bullets to
    state the NEW operating defaults for the upcoming re-run:
      * TSP_CYCLE_LENGTH_OVERRIDE_S: 135s (junction CycleTime) -> 10s
      * REWARD_FUTURE_HORIZON_CYCLES: 1.5 -> 1.0
    including the caveat that _reward_future_term floors cycle_s/horizon_s at
    30s, so the effective lookahead stays 30s rather than literally 10s.

  - New slide "1.4 - Future work (sensitivity sweeps configured, pending run)":
    inserted immediately after slide 1.3. Lists, as a 4-column table:
      * CYCLE_HORIZON_SWEEP (needs re-centring around the new 10s/1-cycle
        baseline -- its existing grid predates that change)
      * DETECTION_SWEEP
      * OCC_SWEEP
      * WOBJ_GAMMA_SWEEP_NASHGATE
      * a new bus-lateness (Z3) vs. NoPriority comparison slide item

  - New slide "2 - ... Preview - target table structure for full results":
    inserted immediately after the "Sensitivity of objectives" slide. Clones
    the 22x8 "Simulation results" KPI table, inserts a new "Bus lateness
    (Z3, Sigma sigma+)" row (populated for ALL strategies incl. NoPriority --
    the new lateness-vs-NoPriority comparison), and replaces all other
    pending data cells with "[TBD]" to show what the table will look like
    once the full re-run (101 configs, new 10s/1-cycle defaults) completes.
"""

import sys
import io
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'

prs = Presentation(PPTX_PATH)


def _set_cell_lines(cell, lines, size_pt=10, bold=False):
    """Set a table cell's text to one paragraph per entry in `lines`,
    applying explicit font size/bold to every run."""
    tf = cell.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = lines[0]
    p0.runs[0].font.size = Pt(size_pt)
    p0.runs[0].font.bold = bold
    for line in lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def _set_cell(cell, text, size_pt=10, bold=False):
    _set_cell_lines(cell, [text], size_pt=size_pt, bold=bold)


# =============================================================================
# Capture clones of slide 7's title/table/footnote BEFORE any slide insertions
# shift indices.
# =============================================================================
s7_orig = prs.slides[7]
shapes7 = {sh.name: sh for sh in s7_orig.shapes}
title7_el_copy = deepcopy(shapes7['title']._element)
table7_el_copy = deepcopy(shapes7['table']._element)
footnote7_el_copy = deepcopy(shapes7['TextBox 14']._element)

# =============================================================================
# STAGE 1: Simplify slide 5 (0-indexed) -- "1.3  Decision-cycle &
# lookahead-horizon configuration"
# =============================================================================
s5 = prs.slides[5]
shapes5 = {sh.name: sh for sh in s5.shapes}

for name in ('Table 202', 'TextBox 203'):
    sh = shapes5[name]
    sh._element.getparent().remove(sh._element)

tb201 = shapes5['TextBox 201']
tf201 = tb201.text_frame
for p in list(tf201.paragraphs[1:]):
    p._p.getparent().remove(p._p)
p0 = tf201.paragraphs[0]
for r in list(p0.runs):
    r._r.getparent().remove(r._r)

bullet_items = [
    ("•  Decision-recalculation cycle (TSP_CYCLE_LENGTH_OVERRIDE_S): ",
     "now defaults to 10s for the upcoming re-run (previously 135s, i.e. each "
     "junction's configured CycleTime). The TSP controller recomputes its "
     "grant/no-action decision for each approaching bus every 10s."),
    ("•  Lookahead horizon (REWARD_FUTURE_HORIZON_CYCLES): ",
     "now defaults to 1.0 × cycle (previously 1.5×). Caveat: _reward_future_term "
     "floors both the effective cycle length and the lookahead horizon at 30s, "
     "so with a 10s cycle the effective lookahead remains 30s ahead rather than "
     "1 × 10s = 10s — revisiting these floors for short decision cycles is "
     "future work (Slide 1.4)."),
    ("•  Sensitivity sweeps over these parameters: ",
     "the existing CYCLE_HORIZON_SWEEP grid (90–180s × 1.0–2.0 cycles) predates "
     "this baseline change and needs re-centring; this, plus the detection-, "
     "occupancy- and weighted-objective (WOBJ_GAMMA) sweeps, is configured but "
     "deferred — see Slide 1.4, Future work."),
]
for i, (lead, rest) in enumerate(bullet_items):
    p = tf201.paragraphs[0] if i == 0 else tf201.add_paragraph()
    r1 = p.add_run()
    r1.text = lead
    r1.font.size = Pt(12)
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(12)
    r2.font.bold = False

tb201.height = Inches(5.0)

# =============================================================================
# STAGE 2: New slide "1.4 - Future work (sensitivity sweeps configured,
# pending run)", inserted immediately after slide 1.3 (index 5 -> 6).
# =============================================================================
shapes4 = {sh.name: sh for sh in prs.slides[4].shapes}

new_slide14 = prs.slides.add_slide(prs.slide_layouts[8])  # "3_Standard slide"
for shape in list(new_slide14.shapes):
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
        shape._element.getparent().remove(shape._element)

title4 = shapes4['TextBox 1']
new_title_el = deepcopy(title4._element)
new_slide14.shapes._spTree.append(new_title_el)
new_title14 = new_slide14.shapes[-1]
cNvPr = new_title14._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '301')
new_title14.text_frame.paragraphs[0].runs[0].text = (
    "1.4  ·  Future work (sensitivity sweeps configured, pending run)"
)

tbl14_shape = new_slide14.shapes.add_table(6, 4, Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.5))
tbl14 = tbl14_shape.table
for ci, w in enumerate([2.6, 5.93, 1.9, 1.7]):
    tbl14.columns[ci].width = Inches(w)
for ri, h in enumerate([0.5, 1.0, 1.0, 1.0, 1.0, 1.0]):
    tbl14.rows[ri].height = Inches(h)

headers14 = ['Sensitivity sweep / item', 'Description', 'Target', 'Status']
for ci, h in enumerate(headers14):
    _set_cell(tbl14.cell(0, ci), h, size_pt=11, bold=True)

rows14 = [
    ("Cycle-length / lookahead-horizon sweep\n(CYCLE_HORIZON_SWEEP)",
     "Re-centre the existing 90–180s × 1.0–2.0-cycle grid around the new "
     "10s / 1-cycle baseline defaults (Slide 1.3); also revisit the 30s "
     "floors in _reward_future_term for short decision cycles.",
     "NashGate (BARGAIN_SPM)",
     "Configured (legacy grid); re-centring pending design"),
    ("Detection-reliability sweep\n(DETECTION_SWEEP)",
     "DETECTION_PROB ∈ {1.00, 0.75, 0.50, 0.25} — probability a ~5s AVL/GPS "
     "bus-tracking cycle registers an in-range bus (1.00 = perfect detection, "
     "current default).",
     "NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
    ("Occupancy-level sweep\n(OCC_SWEEP)",
     "Paired bus/car occupancy (pax/veh) via BUS_OCC_OVERRIDE / "
     "CAR_OCC_OVERRIDE — Low: 20/1.0, Base: 40/1.2, High: 60/1.5.",
     "NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
    ("Weighted-objective α/β/γ sweep\n(WOBJ_GAMMA_SWEEP_NASHGATE)",
     "J = WOBJ_ALPHA·Z1 + WOBJ_BETA·Z2 + WOBJ_GAMMA·Z3 (Z3 = bus lateness, "
     "Σσ⁺); 27 configurations (α∈{0.7,0.8,0.9} × β∈{0.1,0.2,0.3} × "
     "γ∈{0.0,0.1,0.2}).",
     "NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
    ("Bus-lateness vs. NoPriority comparison\n(new results slide)",
     "New Z3 (Σσ⁺, bus lateness) KPI row comparing all 6 strategies "
     "including NoPriority, once Z3 is computed for the no-TSP baseline "
     "(see Preview slide).",
     "All 6 strategies (incl. NoPriority)",
     "Pending full re-run"),
]
for ri, row in enumerate(rows14, start=1):
    for ci, text in enumerate(row):
        _set_cell_lines(tbl14.cell(ri, ci), text.split("\n"), size_pt=10, bold=False)

sldIdLst = prs.slides._sldIdLst
new_sldId14 = sldIdLst[-1]
sldIdLst.remove(new_sldId14)
sldIdLst.insert(6, new_sldId14)

# =============================================================================
# STAGE 3: New "2 - ... Preview - target table structure for full results"
# slide, inserted immediately after "Sensitivity of objectives" (which, after
# Stage 2's insertion, is now at index 9 -> new slide goes to index 10).
# =============================================================================
new_slide_prev = prs.slides.add_slide(prs.slide_layouts[8])
for shape in list(new_slide_prev.shapes):
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
        shape._element.getparent().remove(shape._element)

# -- Title (cloned from slide 7's title) --
new_slide_prev.shapes._spTree.append(title7_el_copy)
new_title_prev = new_slide_prev.shapes[-1]
cNvPr = new_title_prev._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '1001')
runs_t = new_title_prev.text_frame.paragraphs[0].runs
runs_t[2].text = (
    "Preview — target table structure for full results "
    "(101 configs, post re-run with new 10s / 1-cycle defaults)"
)
for r in runs_t:
    r.font.size = Pt(18)

# -- Table (cloned from slide 7's 22x8 KPI table) --
new_slide_prev.shapes._spTree.append(table7_el_copy)
new_table_prev = new_slide_prev.shapes[-1]
cNvPr = new_table_prev._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '1002')
tbl_prev = new_table_prev.table

# Insert a new "Bus lateness (Z3, Sigma sigma+)" row between Z2 (row 19) and
# Z4 (row 20), cloning row 20 as the XML template.
_tbl_prev = tbl_prev._tbl
tr_lst = _tbl_prev.tr_lst
new_tr = deepcopy(tr_lst[20])
tr_lst[20].addprevious(new_tr)
tbl_prev = new_table_prev.table  # re-fetch after XML mutation
n_rows_new = len(tbl_prev.rows)
assert n_rows_new == 23

total_h = Inches(6.21)
for row in tbl_prev.rows:
    row.height = int(total_h / n_rows_new)

TBD = "[TBD]"
DASH_ROWS = {9, 17}              # Side-street avg delay / Throughput-side: leave "--" as-is
NOPRIORITY_DASH_ROWS = {18, 19, 21, 22}  # Z1, Z2, Z4, Total weighted objective: NoPriority stays "--"
Z3_ROW = 20

for ri in range(1, n_rows_new):
    if ri == Z3_ROW:
        _set_cell(tbl_prev.cell(ri, 0), "Bus lateness (Z3, Σσ⁺)", size_pt=9, bold=False)
        _set_cell(tbl_prev.cell(ri, 1), "[bus-s]", size_pt=9, bold=False)
        for ci in range(2, 8):
            _set_cell(tbl_prev.cell(ri, ci), TBD, size_pt=9, bold=False)
        continue
    if ri in DASH_ROWS:
        continue
    for ci in range(2, 8):
        if ri in NOPRIORITY_DASH_ROWS and ci == 2:
            continue
        _set_cell(tbl_prev.cell(ri, ci), TBD, size_pt=9, bold=False)

# -- Footnote (cloned from slide 7's footnote textbox) --
new_slide_prev.shapes._spTree.append(footnote7_el_copy)
new_footnote_prev = new_slide_prev.shapes[-1]
cNvPr = new_footnote_prev._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '1003')
tf_fn = new_footnote_prev.text_frame
for p in list(tf_fn.paragraphs[1:]):
    p._p.getparent().remove(p._p)
p0 = tf_fn.paragraphs[0]
for r in list(p0.runs[1:]):
    r._r.getparent().remove(r._r)
p0.runs[0].text = (
    "[TBD] = to be populated once the full re-run (101 configurations; 10s "
    "decision-recalculation cycle / 1-cycle lookahead defaults, Slide 1.3) "
    "completes. New row ‘Bus lateness (Z3, Σσ⁺)’ is "
    "populated for ALL strategies including NoPriority — enabling a direct "
    "bus-lateness comparison against the no-TSP baseline (Slide 1.4, Future "
    "work). Rows marked ‘—’ (side-street split, throughput-side, "
    "and NoPriority's Z1/Z2/Z4/objective cells) remain not applicable / not "
    "computable from current outputs, as on the preceding slides."
)

sldIdLst = prs.slides._sldIdLst
new_sldId_prev = sldIdLst[-1]
sldIdLst.remove(new_sldId_prev)
sldIdLst.insert(10, new_sldId_prev)

# =============================================================================
# Save
# =============================================================================
prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
print(f"Total slides: {len(prs.slides)}")
