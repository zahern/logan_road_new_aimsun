"""One-off follow-up script (run after update_progress_deck3.py):

  - Slides "2 - ... Simulation results" and "2 - ... Sensitivity of objectives"
    (the 22-row KPI tables): shorten the long title text and reduce its font
    size from 20pt to 18pt so it wraps to 2 lines instead of 3 -- the 3rd line
    was overlapping the table's header row.
"""

from pptx import Presentation
from pptx.util import Pt

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'

prs = Presentation(PPTX_PATH)

# Slide 7 (0-indexed) -- "Simulation results"
s7 = prs.slides[7]
title7 = next(sh for sh in s7.shapes if sh.name == 'title')
runs7 = title7.text_frame.paragraphs[0].runs
runs7[2].text = (
    "Simulation results — Kelvin Grove Rd, AM peak: NoPriority vs. "
    "5 DCTSP strategies (6 of 58 configs; pre-redesign dataset)"
)
for r in runs7:
    r.font.size = Pt(18)

# Slide 8 (0-indexed) -- "Sensitivity of objectives"
s8 = prs.slides[8]
title8 = next(sh for sh in s8.shapes if sh.name == 'title')
runs8 = title8.text_frame.paragraphs[0].runs
runs8[2].text = (
    "Sensitivity of objectives — WOBJ_ALPHA × WOBJ_BETA sweep "
    "(GLOBAL_REWARD+KALMAN; 9 of 58 configs)"
)
for r in runs8:
    r.font.size = Pt(18)

prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
