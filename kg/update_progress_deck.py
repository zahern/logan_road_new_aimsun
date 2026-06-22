"""One-off script: update BCC_progress_meeting_todo.pptx for this fortnight's
progress meeting --

  - Slide 4 ("1.1 - Experiment (scenarios) status"): update status cells for
    the coordinated-TSP and new sensitivity-sweep rows (Kelvin Grove Rd).
  - New slide "1.2 - Sensitivity sweeps configured this period": what the 3
    new sensitivity sweeps (detection / occupancy / demand) test.
  - Slide 5 ("2 - Project Update and Findings: ..."): this period's findings
    (reward-fairness fix, Z2 fix, new sensitivity sweeps) + a chart.
  - Slide 6 ("2 - Project Update and Findings: Simulation results ..."):
    populate the 16-row KPI table for NoPriority + the 5 named DCTSP
    strategies, expanded from 6 to 8 columns.

Source data: batch_results.csv (pre reward-redesign dataset).
"""

from copy import deepcopy

import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH = 'batch_results.csv'
CHART_PATH = '_fig_throughput_vs_delay.png'

# =============================================================================
# 1. Compute slide-6 KPI table values from batch_results.csv
# =============================================================================
STRATS = [
    ("NoPriority", "NO_TSP"),
    ("CPD-QL",     "DCTSP_MARL"),
    ("WaveGate",   "DCTSP_ZIG"),
    ("NashGate",   "DCTSP_BARGAIN_SPM"),
    ("CellSearch", "DCTSP_MP_ECTM"),
    ("CellQLearn", "DCTSP_BXT"),
]
RHO_BUS, RHO_CAR = 40.0, 1.5

df = pd.read_csv(CSV_PATH)
rdata = {}
for label, exp in STRATS:
    rdata[label] = df[df['run_experiment'] == exp].iloc[0]


def fmt(v, nd=1):
    return f"{v:,.{nd}f}"


kpi_rows = []  # (label, unit, [6 values])

# 1. Total bus travel time [min/bus]
kpi_rows.append(("Total bus travel time", "[min/bus]",
                  [fmt(r["stats_Net_TotalTT_h_Bus"] * 60.0 / r["stats_N_DistinctBuses"], 2)
                   for _, r in rdata.items()]))

# 2. Total car travel time [min/car]
kpi_rows.append(("Total car travel time", "[min/car]",
                  [fmt(r["stats_Net_TotalTT_h_Car"] * 60.0 / r["stats_N_DistinctCars"], 2)
                   for _, r in rdata.items()]))

# 3. Total passenger travel time [min/pax]
def pax_tt(r):
    num = (r["stats_Net_TotalTT_h_Car"] * RHO_CAR
           + r["stats_Net_TotalTT_h_Bus"] * RHO_BUS
           + r["stats_Net_TotalTT_h_Truck"] * RHO_CAR) * 60.0
    den = (r["stats_N_DistinctCars"] * RHO_CAR
           + r["stats_N_DistinctBuses"] * RHO_BUS
           + r["stats_N_DistinctTrucks"] * RHO_CAR)
    return num / den

kpi_rows.append(("Total passenger travel time", "[min/pax]",
                  [fmt(pax_tt(r), 2) for _, r in rdata.items()]))

# 4. Total passenger delay [pax-h]
kpi_rows.append(("Total passenger delay", "[pax-h]",
                  [fmt(r["stats_TotalPassDelay_hrs"], 1) for _, r in rdata.items()]))

# 5. Total serviced passengers [pax]
kpi_rows.append(("Total serviced passengers", "[pax]",
                  [fmt(r["stats_PaxEquivPassages"], 0) for _, r in rdata.items()]))

# 6. Average passenger delay [sec/pax]
kpi_rows.append(("Average passenger delay", "[sec/pax]",
                  [fmt(r["stats_AvgPassDelay_s"], 1) for _, r in rdata.items()]))

# 7. Side-street total passenger delay [pax-h]
kpi_rows.append(("Side-street total passenger delay", "[pax-h]",
                  [fmt(r["stats_SidePassDelay_hrs"], 1) for _, r in rdata.items()]))

# 8. Side-street average passenger delay [sec/pax] -- not computable
kpi_rows.append(("Side-street average passenger delay", "[sec/pax]",
                  ["—"] * 6))

# 9. Total system delay [veh-h]
def system_delay_veh_h(r):
    total_veh = r["stats_N_DistinctCars"] + r["stats_N_DistinctBuses"] + r["stats_N_DistinctTrucks"]
    return r["stats_Net_Delay_All"] * total_veh / 3600.0

kpi_rows.append(("Total system delay", "[veh-h]",
                  [fmt(system_delay_veh_h(r), 1) for _, r in rdata.items()]))

# 10. Total hours of travel [veh-h]
kpi_rows.append(("Total hours of travel", "[veh-h]",
                  [fmt(r["stats_Net_TotalTT_h_Car"] + r["stats_Net_TotalTT_h_Bus"] + r["stats_Net_TotalTT_h_Truck"], 1)
                   for _, r in rdata.items()]))

# 11. Total VKT - car [veh-km]
kpi_rows.append(("Total VKT – car", "[veh-km]",
                  [fmt(r["stats_Net_TotalDist_Car"], 0) for _, r in rdata.items()]))

# 12. Total VKT - bus [veh-km]
kpi_rows.append(("Total VKT – bus", "[veh-km]",
                  [fmt(r["stats_Net_TotalDist_Bus"], 2) for _, r in rdata.items()]))

# 13. Total vehicles in system [veh]
def total_veh(r):
    return r["stats_N_DistinctCars"] + r["stats_N_DistinctBuses"] + r["stats_N_DistinctTrucks"]

kpi_rows.append(("Total vehicles in system", "[veh]",
                  [fmt(total_veh(r), 0) for _, r in rdata.items()]))

# 14. Total vehicles in system (Bus) [veh]
kpi_rows.append(("Total vehicles in system (Bus)", "[veh]",
                  [fmt(r["stats_N_DistinctBuses"], 0) for _, r in rdata.items()]))

# 15. Throughput - main [veh/h]
kpi_rows.append(("Throughput – main", "[veh/h]",
                  [fmt(r["stats_Net_Flow_Car"] + r["stats_Net_Flow_Bus"] + r["stats_Net_Flow_Truck"], 1)
                   for _, r in rdata.items()]))

# 16. Throughput - side [veh/h] -- not computable
kpi_rows.append(("Throughput – side", "[veh/h]", ["—"] * 6))

assert len(kpi_rows) == 16

# =============================================================================
# 2. Chart for slide 5: total vehicles served vs avg passenger delay
# =============================================================================
total_veh_vals = [total_veh(r) for _, r in rdata.items()]
avg_delay_vals = [r["stats_AvgPassDelay_s"] for _, r in rdata.items()]
labels = [label for label, _ in STRATS]
corr_r = pd.Series(total_veh_vals).corr(pd.Series(avg_delay_vals))

fig, ax = plt.subplots(figsize=(5.89, 3.17), dpi=150)
ax.scatter(total_veh_vals, avg_delay_vals, s=60, color="#00467F", zorder=3)
for x, y, lab in zip(total_veh_vals, avg_delay_vals, labels):
    ax.annotate(lab, (x, y), textcoords="offset points", xytext=(6, 4), fontsize=8)
ax.set_xlabel("Total vehicles served in network [veh]", fontsize=9)
ax.set_ylabel("Average passenger delay [s/pax]", fontsize=9)
ax.set_title(f"Throughput vs. passenger delay  (r = {corr_r:.2f})", fontsize=10)
ax.tick_params(labelsize=8)
ax.grid(True, linestyle=":", linewidth=0.5, alpha=0.6)
fig.tight_layout()
fig.savefig(CHART_PATH)
plt.close(fig)

# =============================================================================
# 3. Open presentation
# =============================================================================
prs = Presentation(PPTX_PATH)


def _replace_cell_text(cell, text):
    """Replace a table cell's text in-place, preserving the first run's
    formatting (font size/bold/color) and dropping any extra runs/paragraphs."""
    tf = cell.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)
    if p.runs:
        p.runs[0].text = text
    else:
        run = p.add_run()
        run.text = text


def _set_cell(cell, text, size_pt=9, bold=False):
    """Replace a table cell's text and set explicit font size/bold."""
    _replace_cell_text(cell, text)
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(size_pt)
    run.font.bold = bold


# =============================================================================
# B1. Slide 4 -- update experiment-status table cells
# =============================================================================
s4 = prs.slides[3]
tbl4 = next(sh for sh in s4.shapes if sh.has_table).table

# Row indices (0 = header):
#  3 = Coordinated TSP - Kelvin Grove Rd
#  5 = Sensitivity - demand levels: Kelvin Grove Rd
#  7 = Sensitivity - occupancy levels: Kelvin Grove Rd
#  9 = Sensitivity - detection levels: Kelvin Grove Rd
_replace_cell_text(tbl4.cell(3, 1), "Done")
_replace_cell_text(tbl4.cell(3, 2), "Done")
_replace_cell_text(tbl4.cell(3, 3), "Done")
_replace_cell_text(tbl4.cell(3, 4), "In development")

for _row in (5, 7, 9):
    _replace_cell_text(tbl4.cell(_row, 1), "Done")
    _replace_cell_text(tbl4.cell(_row, 2), "In development")
    _replace_cell_text(tbl4.cell(_row, 3), "Pending")
    _replace_cell_text(tbl4.cell(_row, 4), "–")

s4_title = next(sh for sh in s4.shapes if sh.name == "title")

# Capture references to slides 5 and 6 BEFORE B2 reorders _sldIdLst (which
# would otherwise shift prs.slides[4]/[5] to point at the wrong slides).
s5 = prs.slides[4]
s6 = prs.slides[5]

# =============================================================================
# B2. New slide "1.2 - Sensitivity sweeps configured this period"
# =============================================================================
new_slide = prs.slides.add_slide(prs.slide_layouts[8])  # "3_Standard slide" (same as slide 4/5/6)

# Remove the auto-added empty Title placeholder (slides 4-6 use a custom
# "title" textbox instead, not the layout's title placeholder).
for shape in list(new_slide.shapes):
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
        shape._element.getparent().remove(shape._element)

# Title textbox -- match slide 4's "1.1 -" title style exactly.
title_box = new_slide.shapes.add_textbox(Inches(0.6), Inches(0.328), Inches(12.133), Inches(0.897))
tf = title_box.text_frame
run = tf.paragraphs[0].add_run()
run.text = "1.2  ·  Sensitivity sweeps configured this period"
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = s4_title.text_frame.paragraphs[0].runs[0].font.color.rgb

# Table: 4 cols x 4 rows.
n_rows, n_cols = 4, 4
gtable = new_slide.shapes.add_table(n_rows, n_cols, Inches(0.605), Inches(1.5), Inches(12.133), Inches(4.0)).table
col_widths = [1.8, 6.0, 1.6, 2.733]
for ci, w in enumerate(col_widths):
    gtable.columns[ci].width = Inches(w)
gtable.rows[0].height = Inches(0.45)
for ri in range(1, n_rows):
    gtable.rows[ri].height = Inches(1.18)

header = ["Sensitivity sweep", "Parameter(s) & levels", "Target", "Status"]
for ci, text in enumerate(header):
    _set_cell(gtable.cell(0, ci), text, size_pt=10.5, bold=True)

sweep_rows = [
    ("Detection levels",
     "DETECTION_PROB ∈ {1.00, 0.75, 0.50, 0.25} — probability that a ~5s "
     "AVL/GPS bus-tracking cycle registers an in-range bus (1.00 = perfect "
     "detection, current default)",
     "NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
    ("Occupancy levels",
     "Paired bus/car occupancy assumptions (pax/veh) via BUS_OCC_OVERRIDE / "
     "CAR_OCC_OVERRIDE — Low: 20/1.0, Base: 40/1.2, High: 60/1.5 "
     "(car value also applied to trucks)",
     "NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
    ("Demand levels",
     "Demand scalar ∈ {0.8, 1.0, 1.2} applied network-wide to all "
     "traffic-demand matrices (±20% of AM/PM demand)",
     "NoPriority + NashGate (BARGAIN_SPM)",
     "Configured, pending run"),
]
for ri, (sweep, params, target, status) in enumerate(sweep_rows, start=1):
    _set_cell(gtable.cell(ri, 0), sweep, size_pt=10, bold=False)
    _set_cell(gtable.cell(ri, 1), params, size_pt=10, bold=False)
    _set_cell(gtable.cell(ri, 2), target, size_pt=10, bold=False)
    _set_cell(gtable.cell(ri, 3), status, size_pt=10, bold=False)

# Reposition the new slide immediately after slide 4 (index 3 -> position 4).
sldIdLst = prs.slides._sldIdLst
new_sldId = sldIdLst[-1]
sldIdLst.remove(new_sldId)
sldIdLst.insert(4, new_sldId)

# =============================================================================
# B3. Slide 5 -- this period's findings
# =============================================================================
shapes5 = {sh.name: sh for sh in s5.shapes}

# Title
title_runs = shapes5["title"].text_frame.paragraphs[0].runs
title_runs[2].text = "Reward-function fairness fix + sensitivity-sweep scope"

# Body bullets
body_tf = shapes5["body"].text_frame
p0, p1, p2 = body_tf.paragraphs[0], body_tf.paragraphs[1], body_tf.paragraphs[2]
for p in (p0, p1, p2):
    for r in list(p.runs[1:]):
        r._r.getparent().remove(r._r)

p0.runs[0].text = (
    f"Every coordinated TSP strategy serves fewer total vehicles than "
    f"NoPriority ({min(total_veh_vals):,.0f}–{max(v for v in total_veh_vals if v != max(total_veh_vals)):,.0f} "
    f"vs {max(total_veh_vals):,.0f}); across the 6 strategies this anti-correlates "
    f"with average passenger delay (r = {corr_r:.2f}) — part of each strategy's "
    f"passenger-delay ‘win’ reflects reduced network throughput, not just "
    f"better bus priority."
)
p1.runs[0].text = (
    "– Fixed by adding a dynamic network-throughput factor (NETWORK_FACTOR, "
    "ramped 1.0→5.0 with live network density) to the shared cross-traffic "
    "cost term used by all 6 TSP reward modes, and fixed the Z2 "
    "(offset-correction) dead-code path — was always 0 due to a "
    "‘.gb is None’ filter mismatch."
)
p2.runs[0].text = (
    "– Added 3 new sensitivity-sweep mechanisms (detection reliability, "
    "bus/car occupancy, demand level), configured against NashGate — the "
    "current lowest-delay strategy. See slide 1.2. All pending the full "
    "re-run with the fixes above."
)

# Replace the "fig" placeholder autoshape with the generated chart.
fig_shape = shapes5["fig"]
left, top, width, height = fig_shape.left, fig_shape.top, fig_shape.width, fig_shape.height
fig_shape._element.getparent().remove(fig_shape._element)
s5.shapes.add_picture(CHART_PATH, left, top, width=width, height=height)

# =============================================================================
# B4. Slide 6 -- populate the 16-row KPI table (expand 6 -> 8 columns)
# =============================================================================
shapes6 = {sh.name: sh for sh in s6.shapes}

# Title
title_runs6 = shapes6["title"].text_frame.paragraphs[0].runs
title_runs6[2].text = (
    "Simulation results — Kelvin Grove Rd corridor, AM peak: "
    "NoPriority vs. 5 DCTSP strategies (current dataset — pre reward-redesign re-run)"
)

tbl6_shape = shapes6["table"]
tbl6 = tbl6_shape.table
_tbl = tbl6._tbl
tblGrid = _tbl.tblGrid

# Add 2 columns (clone the last gridCol) and 2 cells per row (clone the last
# tc, which is empty for data rows / "..." for the header row).
last_gridCol = tblGrid.gridCol_lst[-1]
for _ in range(2):
    tblGrid.append(deepcopy(last_gridCol))

for tr in _tbl.tr_lst:
    tcs = tr.tc_lst
    last_tc = tcs[-1]
    for _ in range(2):
        tr.append(deepcopy(last_tc))

# Re-fetch the table view now that the XML has more columns/cells.
tbl6 = tbl6_shape.table

# Column widths: KPI | unit | 6 data cols, total ~11.6"
col_widths = [2.7, 0.85] + [1.3417] * 6
for ci, w in enumerate(col_widths):
    tbl6.columns[ci].width = Inches(w)

# Header row
header_cols = ["KPI", "unit"] + [label for label, _ in STRATS]
for ci, text in enumerate(header_cols):
    _set_cell(tbl6.cell(0, ci), text, size_pt=9, bold=True)

# KPI rows
for ri, (label, unit, values) in enumerate(kpi_rows, start=1):
    _set_cell(tbl6.cell(ri, 0), label, size_pt=9, bold=False)
    _set_cell(tbl6.cell(ri, 1), unit, size_pt=9, bold=False)
    for ci, val in enumerate(values, start=2):
        _set_cell(tbl6.cell(ri, ci), val, size_pt=9, bold=False)

# =============================================================================
# Save
# =============================================================================
prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
print(f"Chart: {CHART_PATH} (r = {corr_r:.3f})")
print(f"Total slides: {len(prs.slides)}")
