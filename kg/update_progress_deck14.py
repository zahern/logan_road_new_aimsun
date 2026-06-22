"""
update_progress_deck14.py — Populate bcc_progress_meeting_16jun_fake.pptx with
fake/illustrative data for NO_TSP + WaveGate (DCTSP_ZIG) sensitivity sweeps.

Slide 9  (idx 8): Occupancy sweep — NoPriority | WaveGate | OCC_LOW | OCC_BASE | OCC_HIGH
Slide 10 (idx 9): Detection sweep — NoPriority ref | WaveGate P=1.00 | P=0.75 | P=0.50 | P=0.25
Slide 11 (idx10): Demand sweep   — NO_TSP×0.8/1.0/1.2 | WaveGate×0.8/1.0/1.2

NO_TSP and WaveGate base: REAL values from batch_results.csv.
Sweep variants: synthesised from the two bases at plausible offsets.

After the Aimsun batch runs (batch_runner.py config: NO_TSP + DCTSP_ZIG +
OCC_SWEEP_WAVEGATE + DETECTION_SWEEP_WAVEGATE + DEMAND_SWEEP), re-run with the
real CSV to replace these fake values.
"""

import math, shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

SRC_PPTX = 'BCC_progress_meeting_todo.pptx'
OUT_PPTX = 'bcc_progress_meeting_16jun_fake.pptx'
CSV_PATH = 'batch_results.csv'

RHO_BUS, RHO_CAR = 40.0, 1.5

# ── Formatting ────────────────────────────────────────────────────────────────
def f1(x):  return f'{x:,.1f}'
def f2(x):  return f'{x:,.2f}'
def fi(x):  return f'{int(round(x)):,}'
FMT = {
    1:f1, 2:f2, 3:f2, 4:f2, 5:f1, 6:fi, 7:f1, 8:f1,
    10:f1, 11:f1, 12:fi, 13:f2, 14:fi, 15:fi, 16:f1,
    18:f1, 19:f1, 20:f1, 21:f1, 22:f1,
}

# ── KPI formula ───────────────────────────────────────────────────────────────
def compute_kpi(d):
    r = {}
    r[1]  = d['TT_Bus']
    r[2]  = d['TT_Bus'] * 60.0 / d['N_Bus']
    r[3]  = d['TT_Car'] * 60.0 / d['N_Car']
    num4  = (d['TT_Car']*RHO_CAR + d['TT_Bus']*RHO_BUS + d['TT_Trk']*RHO_CAR) * 60.0
    den4  = d['N_Car']*RHO_CAR + d['N_Bus']*RHO_BUS + d['N_Trk']*RHO_CAR
    r[4]  = num4 / den4
    r[5]  = d['PassDelay']
    r[6]  = d['PaxPass']
    r[7]  = d['AvgPassDelay']
    r[8]  = d['SideDelay']
    totveh = d['N_Car'] + d['N_Bus'] + d['N_Trk']
    r[10] = d['NetDelay'] * totveh / 3600.0
    r[11] = d['TT_Car'] + d['TT_Bus'] + d['TT_Trk']
    r[12] = d['DistCar']
    r[13] = d['DistBus']
    r[14] = totveh
    r[15] = d['N_Bus']
    r[16] = d['FlowCar'] + d['FlowBus'] + d['FlowTrk']
    r[18] = d['Z1']
    r[19] = d['Z2']
    r[20] = d['Z3']
    r[21] = d['Z4']
    r[22] = d['Obj']
    return r

# ── Load real NO_TSP / WaveGate values ────────────────────────────────────────
_raw = pd.read_csv(CSV_PATH)

def _row(name):
    r = _raw[_raw['run_experiment'] == name].iloc[0]
    return {
        'TT_Bus':    r['stats_Net_TotalTT_h_Bus'],
        'N_Bus':     r['stats_N_DistinctBuses'],
        'TT_Car':    r['stats_Net_TotalTT_h_Car'],
        'N_Car':     r['stats_N_DistinctCars'],
        'PassDelay': r['stats_TotalPassDelay_hrs'],
        'PaxPass':   r['stats_PaxEquivPassages'],
        'AvgPassDelay': r['stats_AvgPassDelay_s'],
        'SideDelay': r['stats_SidePassDelay_hrs'],
        'NetDelay':  r['stats_Net_Delay_All'],
        'TT_Trk':    r['stats_Net_TotalTT_h_Truck'],
        'N_Trk':     r['stats_N_DistinctTrucks'],
        'DistCar':   r['stats_Net_TotalDist_Car'],
        'DistBus':   r['stats_Net_TotalDist_Bus'],
        'FlowCar':   r['stats_Net_Flow_Car'],
        'FlowBus':   r['stats_Net_Flow_Bus'],
        'FlowTrk':   r['stats_Net_Flow_Truck'],
        'Z1':  r['wobj_Z1_total'],
        'Z2':  r['wobj_Z2_total'],
        'Z3':  r['wobj_Z3_total'],
        'Z4':  r['wobj_Z4_total'],
        'Obj': r['wobj_objective_total'],
    }

NO_TSP   = _row('NO_TSP')
NASHGATE = _row('DCTSP_BARGAIN_SPM')   # real data we have

# WaveGate (DCTSP_ZIG) is not in the CSV yet; synthesise from NashGate ×1.08 offset
# (consistent with the prior 6-strategy fake — slightly worse Z1 than NashGate).
_NG = NASHGATE
WG_BASE = {
    'TT_Bus':    _NG['TT_Bus']    * 1.040,  'N_Bus':  int(round(_NG['N_Bus']  * 1.000)),
    'TT_Car':    _NG['TT_Car']    * 1.050,  'N_Car':  int(round(_NG['N_Car']  * 0.988)),
    'PassDelay': _NG['PassDelay'] * 1.069,  'PaxPass': int(round(_NG['PaxPass'] * 0.990)),
    'AvgPassDelay': _NG['AvgPassDelay'] * 1.059,
    'SideDelay': _NG['SideDelay'] * 1.069,
    'NetDelay':  _NG['NetDelay']  * 1.070,
    'TT_Trk':   _NG['TT_Trk']   * 1.029,  'N_Trk': int(round(_NG['N_Trk']  * 1.005)),
    'DistCar':   _NG['DistCar']  * 0.990,  'DistBus': _NG['DistBus'] * 0.980,
    'FlowCar':   _NG['FlowCar']  * 0.980,  'FlowBus': _NG['FlowBus'] * 0.970,
    'FlowTrk':   _NG['FlowTrk']  * 0.969,
    'Z1':  _NG['Z1']  * 1.080,  'Z2':  _NG['Z2']  * 0.920,
    'Z3':  _NG['Z3']  * 1.020,  'Z4':  _NG['Z4']  * 1.050,
    'Obj': _NG['Obj'] * 1.080,
}

# ── Synthesised sweep data ────────────────────────────────────────────────────
def _scale(base, **mults):
    """Return a copy of base dict with each key scaled by its multiplier."""
    d = dict(base)
    for k, m in mults.items():
        if k in d:
            d[k] = d[k] * m
    return d

# --- Occupancy sweep (bus_occ / car_occ): WaveGate only ---
# OCC_LOW  (bus=20, car=1.0): TSP less aggressive → worse Z1 & objective
# OCC_BASE (bus=40, car=1.2): cross-check of WaveGate base (slight noise)
# OCC_HIGH (bus=60, car=1.5): TSP more aggressive → better Z1 & objective
OCC_LOW  = _scale(WG_BASE, TT_Bus=1.12, TT_Car=0.95, N_Car=0.98,
                  PassDelay=1.22, PaxPass=0.97, AvgPassDelay=1.18,
                  SideDelay=1.20, NetDelay=1.15,
                  TT_Trk=1.08, DistCar=0.97, DistBus=0.96,
                  FlowCar=0.95, FlowBus=0.94, FlowTrk=0.95,
                  Z1=1.25, Z2=0.85, Z3=1.15, Z4=1.12, Obj=1.25)

OCC_BASE = _scale(WG_BASE, TT_Bus=1.02, TT_Car=1.01, PassDelay=1.02,
                  PaxPass=0.99, AvgPassDelay=1.02, SideDelay=1.01,
                  Z1=1.02, Z2=0.99, Z3=1.01, Z4=1.01, Obj=1.02)

OCC_HIGH = _scale(WG_BASE, TT_Bus=0.93, TT_Car=1.08, N_Car=1.01,
                  PassDelay=0.88, PaxPass=1.01, AvgPassDelay=0.87,
                  SideDelay=0.87, NetDelay=1.05,
                  TT_Trk=0.97, DistCar=1.01, DistBus=1.01,
                  FlowCar=1.02, FlowBus=1.02, FlowTrk=1.01,
                  Z1=0.85, Z2=1.12, Z3=0.93, Z4=0.97, Obj=0.85)

# --- Detection sweep (DETECTION_PROB): WaveGate only ---
# P100: same as WaveGate base (perfect detection)
# P075: 25% missed → TSP fires 25% less often → 18% worse Z1
# P050: 50% missed → 40% worse Z1
# P025: 75% missed → performance collapses toward NoPriority
DET_P100 = dict(WG_BASE)
DET_P075 = _scale(WG_BASE, TT_Bus=1.10, TT_Car=0.96, PassDelay=1.15,
                  PaxPass=0.98, AvgPassDelay=1.12, SideDelay=1.08,
                  NetDelay=0.96, TT_Trk=0.97, DistCar=0.98, FlowCar=0.98,
                  Z1=1.18, Z2=0.88, Z3=1.12, Z4=0.98, Obj=1.18)
DET_P050 = _scale(WG_BASE, TT_Bus=1.25, TT_Car=0.92, PassDelay=1.35,
                  PaxPass=0.96, AvgPassDelay=1.30, SideDelay=1.20,
                  NetDelay=0.90, TT_Trk=0.93, DistCar=0.96, FlowCar=0.96,
                  Z1=1.42, Z2=0.72, Z3=1.28, Z4=0.93, Obj=1.42)
DET_P025 = _scale(WG_BASE, TT_Bus=1.60, TT_Car=0.82, PassDelay=1.75,
                  PaxPass=0.94, AvgPassDelay=1.68, SideDelay=1.45,
                  NetDelay=0.75, TT_Trk=0.83, DistCar=0.94, FlowCar=0.92,
                  Z1=2.10, Z2=0.40, Z3=1.58, Z4=0.83, Obj=2.10)

# --- Demand sweep (demand scalar 0.8 / 1.0 / 1.2) for NO_TSP and WaveGate ---
# Lower demand: less congestion (sub-linear TT reduction), fewer vehicles, lower absolute Z
# Higher demand: more congestion (super-linear TT increase), more vehicles
NO_TSP_D08 = _scale(NO_TSP, TT_Bus=0.90, N_Bus=0.95,
                    TT_Car=0.75, N_Car=0.82,
                    PassDelay=0.78, PaxPass=0.82, AvgPassDelay=0.82,
                    SideDelay=0.76, NetDelay=0.72,
                    TT_Trk=0.78, N_Trk=0.82,
                    DistCar=0.82, DistBus=0.88,
                    FlowCar=0.82, FlowBus=0.88, FlowTrk=0.82)
NO_TSP_D10 = dict(NO_TSP)  # base (real values)
NO_TSP_D12 = _scale(NO_TSP, TT_Bus=1.12, N_Bus=1.05,
                    TT_Car=1.38, N_Car=1.18,
                    PassDelay=1.32, PaxPass=1.18, AvgPassDelay=1.22,
                    SideDelay=1.30, NetDelay=1.40,
                    TT_Trk=1.28, N_Trk=1.18,
                    DistCar=1.18, DistBus=1.10,
                    FlowCar=1.18, FlowBus=1.10, FlowTrk=1.18)

WG_D08 = _scale(WG_BASE, TT_Bus=0.92, N_Bus=0.95,
                TT_Car=0.78, N_Car=0.83,
                PassDelay=0.82, PaxPass=0.83, AvgPassDelay=0.83,
                SideDelay=0.78, NetDelay=0.80,
                TT_Trk=0.80, N_Trk=0.82,
                DistCar=0.83, DistBus=0.89,
                FlowCar=0.83, FlowBus=0.89, FlowTrk=0.83,
                Z1=0.83, Z2=0.85, Z3=0.86, Z4=0.79, Obj=0.83)
WG_D10  = dict(WG_BASE)
WG_D12  = _scale(WG_BASE, TT_Bus=1.18, N_Bus=1.05,
                 TT_Car=1.42, N_Car=1.18,
                 PassDelay=1.35, PaxPass=1.18, AvgPassDelay=1.28,
                 SideDelay=1.32, NetDelay=1.45,
                 TT_Trk=1.32, N_Trk=1.18,
                 DistCar=1.18, DistBus=1.10,
                 FlowCar=1.18, FlowBus=1.10, FlowTrk=1.18,
                 Z1=1.30, Z2=1.18, Z3=1.22, Z4=1.42, Obj=1.30)

# ── NaN for NO_TSP Z-objectives (no TSP controller) ───────────────────────────
for _d in (NO_TSP_D08, NO_TSP_D10, NO_TSP_D12):
    _d['Z1'] = float('nan')
    _d['Z2'] = float('nan')
    _d['Z3'] = float('nan')
    _d['Obj'] = float('nan')

# ── XML helpers ───────────────────────────────────────────────────────────────
A_NS         = 'http://schemas.openxmlformats.org/drawingml/2006/main'
HEADER_FILL  = '00467F'
NP_FILL      = 'FFF2CC'
BORDER_COLOR = 'BFBFBF'
BORDER_W     = 6350

def _color_fill(hex_c):
    sf = etree.Element(f'{{{A_NS}}}solidFill')
    etree.SubElement(sf, f'{{{A_NS}}}srgbClr', val=hex_c)
    return sf

def _border(tag):
    ln = etree.Element(f'{{{A_NS}}}{tag}',
                       w=str(BORDER_W), cap='flat', cmpd='sng', algn='ctr')
    ln.append(_color_fill(BORDER_COLOR))
    etree.SubElement(ln, f'{{{A_NS}}}prstDash', val='solid')
    etree.SubElement(ln, f'{{{A_NS}}}round')
    etree.SubElement(ln, f'{{{A_NS}}}headEnd', type='none', w='med', len='med')
    etree.SubElement(ln, f'{{{A_NS}}}tailEnd', type='none', w='med', len='med')
    return ln

def _style(cell, fill_hex=None, no_fill=False, mar_lr=73152, mar_tb=27432):
    for x in cell._tc.findall(f'{{{A_NS}}}tcPr'):
        cell._tc.remove(x)
    tcPr = etree.SubElement(cell._tc, f'{{{A_NS}}}tcPr',
                             marL=str(mar_lr), marR=str(mar_lr),
                             marT=str(mar_tb), marB=str(mar_tb), anchor='ctr')
    for t in ('lnL', 'lnR', 'lnT', 'lnB'):
        tcPr.append(_border(t))
    if fill_hex:
        tcPr.append(_color_fill(fill_hex))
    elif no_fill:
        etree.SubElement(tcPr, f'{{{A_NS}}}noFill')

def _text(cell, text, bold=False, color=None, sz=800, align='ctr', line2=None):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p._p.insert(0, etree.Element(f'{{{A_NS}}}pPr', algn=align))
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz / 100)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor.from_string(color)
    if line2 is not None:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = line2
        r2.font.size = Pt(sz / 100)
        r2.font.bold = bold
        if color:
            r2.font.color.rgb = RGBColor.from_string(color)

# ── Row definitions (same KPI structure across all three slides) ──────────────
ROWS = [
    ('Total bus travel time',               '[bus-h]',       1),
    ('Average bus travel time',             '[min/bus]',     2),
    ('Total car travel time',               '[min/car]',     3),
    ('Total passenger travel time',         '[min/pax]',     4),
    ('Total passenger delay',               '[pax-h]',       5),
    ('Total serviced passengers',           '[pax]',         6),
    ('Average passenger delay',             '[sec/pax]',     7),
    ('Side-street total passenger delay',   '[pax-h]',       8),
    ('Side-street average passenger delay', '[sec/pax]',    'dash'),
    ('Total system delay',                  '[veh-h]',       10),
    ('Total hours of travel',               '[veh-h]',       11),
    ('Total VKT – car',               '[veh-km]',     12),
    ('Total VKT – bus',               '[veh-km]',     13),
    ('Total vehicles in system',            '[veh (count)]', 14),
    ('Total vehicles in system (Bus)',      '[veh (count)]', 15),
    ('Throughput – main',             '[veh/h]',       16),
    ('Throughput – side',             '[veh/h]',      'dash'),
    ('Weighted pax-delay obj. (Z1)',        '[a.u.]',        18),
    ('Offset-correction obj. (Z2)',         '[a.u.]',        19),
    ('Raw travel-time obj. (Z4)',           '[a.u.]',        21),
    ('Total weighted objective',            '[a.u.]',        22),
]

# ── Generic table builder ─────────────────────────────────────────────────────
def build_table(slide, col_specs, col_widths_in, row_h_in, title_text, fn_text):
    """
    col_specs: list of (header_line1, header_line2_or_None, data_dict, is_nopriority)
               data_dict: raw stats dict — passed through compute_kpi()
               is_nopriority: True → yellow fill, Z-rows (18-22) show '—'
    """
    old = next(sh for sh in slide.shapes if sh.has_table)
    L, T, W, H = old.left, old.top, old.width, old.height
    old._element.getparent().remove(old._element)

    NR = 1 + len(ROWS)
    NC = 2 + len(col_specs)
    tsh = slide.shapes.add_table(NR, NC, L, T, W, H)
    tbl = tsh.table

    for ci, w in enumerate(col_widths_in):
        tbl.columns[ci].width = Inches(w)
    rh = Inches(row_h_in)
    for ri in range(NR):
        tbl.rows[ri].height = rh

    # Header row
    _style(tbl.cell(0, 0), fill_hex=HEADER_FILL)
    _text(tbl.cell(0, 0), 'KPI', bold=True, color='FFFFFF', align='l')
    _style(tbl.cell(0, 1), fill_hex=HEADER_FILL)
    _text(tbl.cell(0, 1), 'unit', bold=True, color='FFFFFF')
    for ci, (h1, h2, _, _is_np) in enumerate(col_specs):
        _style(tbl.cell(0, 2 + ci), fill_hex=HEADER_FILL)
        _text(tbl.cell(0, 2 + ci), h1, bold=True, color='FFFFFF', line2=h2)

    # Pre-compute KPI series for each column
    kpi_list = []
    for _, _, ddict, _ in col_specs:
        df1 = pd.DataFrame([ddict])
        kpi_list.append(compute_kpi(df1))

    Z_ROWS = {18, 19, 20, 21, 22}

    for ri, (label, unit, fkey) in enumerate(ROWS):
        trow = ri + 1
        _style(tbl.cell(trow, 0), no_fill=True)
        _text(tbl.cell(trow, 0), label, align='l')
        _style(tbl.cell(trow, 1), no_fill=True)
        _text(tbl.cell(trow, 1), unit)
        for ci, (_, _, _, is_np) in enumerate(col_specs):
            fill = NP_FILL if is_np else None
            _style(tbl.cell(trow, 2 + ci), fill_hex=fill, no_fill=(not is_np))
            if fkey == 'dash':
                txt = '—'
            elif is_np and trow in Z_ROWS:
                txt = '—'
            else:
                val = kpi_list[ci][fkey].iloc[0]
                txt = '—' if math.isnan(val) else FMT[fkey](val)
            _text(tbl.cell(trow, 2 + ci), txt)

    # Update title and footnote
    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table:
            runs = sh.text_frame.paragraphs[0].runs
            if sh.top < 200_000 and len(runs) >= 3:
                runs[2].text = title_text
            elif sh.top > 4_000_000 and runs:
                runs[0].text = fn_text


# ── Build fake PPTX ───────────────────────────────────────────────────────────
shutil.copy(SRC_PPTX, OUT_PPTX)
prs = Presentation(OUT_PPTX)

DRAFT = '(DRAFT — 16 Jun 2026, fake — pending Aimsun re-run)'
FN = (
    f'Fake/illustrative data {DRAFT}. '
    'NO_TSP values are real (batch_results.csv 2026-06-16). '
    'WaveGate (DCTSP_ZIG) base and all sweep variants are synthesised at plausible offsets '
    'from the real NashGate (DCTSP_BARGAIN_SPM) result. '
    'Run batch_runner.py (NO_TSP + DCTSP_ZIG + OCC/DETECTION sweeps enabled) to obtain real data.'
)

# ── Slide 9 — Occupancy sensitivity (WaveGate) ───────────────────────────────
build_table(
    slide=prs.slides[8],
    col_specs=[
        ('NoPriority', None,       NO_TSP,   True),
        ('WaveGate',   'base',     WG_BASE,  False),
        ('OCC Low',    'bus=20\ncar=1.0', OCC_LOW,  False),
        ('OCC Base',   'bus=40\ncar=1.2', OCC_BASE, False),
        ('OCC High',   'bus=60\ncar=1.5', OCC_HIGH, False),
    ],
    col_widths_in=[2.4, 0.55, 1.73, 1.73, 1.73, 1.73, 1.73],
    row_h_in=0.282,
    title_text=(
        'Sensitivity — occupancy levels: NoPriority vs WaveGate (DCTSP_ZIG) '
        f'[bus_occ ∈ {{20, 40, 60}} pax/veh; 10s/1-cycle defaults] {DRAFT}'
    ),
    fn_text=FN,
)

# ── Slide 10 — Detection sensitivity (WaveGate) ──────────────────────────────
build_table(
    slide=prs.slides[9],
    col_specs=[
        ('NoPriority', '(ref)',   NO_TSP,   True),
        ('WaveGate',   'P=1.00', DET_P100, False),
        ('WaveGate',   'P=0.75', DET_P075, False),
        ('WaveGate',   'P=0.50', DET_P050, False),
        ('WaveGate',   'P=0.25', DET_P025, False),
    ],
    col_widths_in=[2.4, 0.55, 1.73, 1.73, 1.73, 1.73, 1.73],
    row_h_in=0.282,
    title_text=(
        'Sensitivity — detection reliability: WaveGate (DCTSP_ZIG) '
        f'[DETECTION_PROB ∈ {{1.00, 0.75, 0.50, 0.25}}; 10s/1-cycle defaults] {DRAFT}'
    ),
    fn_text=FN,
)

# ── Slide 11 — Demand sensitivity (NO_TSP + WaveGate) ────────────────────────
build_table(
    slide=prs.slides[10],
    col_specs=[
        ('NO_TSP',    '×0.8', NO_TSP_D08, True),
        ('NO_TSP',    '×1.0', NO_TSP_D10, True),
        ('NO_TSP',    '×1.2', NO_TSP_D12, True),
        ('WaveGate',  '×0.8', WG_D08,     False),
        ('WaveGate',  '×1.0', WG_D10,     False),
        ('WaveGate',  '×1.2', WG_D12,     False),
    ],
    col_widths_in=[2.7, 0.55, 1.39, 1.39, 1.39, 1.39, 1.39, 1.39],
    row_h_in=0.27,
    title_text=(
        'Sensitivity — demand levels: NoPriority & WaveGate (DCTSP_ZIG) '
        f'[demand scalar ∈ {{0.8, 1.0, 1.2}}; 10s/1-cycle defaults] {DRAFT}'
    ),
    fn_text=FN,
)

prs.save(OUT_PPTX)
print('Saved', OUT_PPTX)
