"""One-off follow-up script (run after update_progress_deck6.py):

  - "Preview" slide: the 23-row table still overflowed the footer band --
    the narrow "unit" column wrapped "[veh (count)]" onto 2 lines, inflating
    those rows. Widen the unit column (0.85in -> 1.05in), shrink the KPI
    column slightly to compensate (2.70in -> 2.60in), shrink body font to
    8pt, and move the footnote up to just below the (now shorter) table.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
prs = Presentation(PPTX_PATH)

s10 = prs.slides[10]

tbl_shape = next(sh for sh in s10.shapes if sh.has_table)
tbl = tbl_shape.table

tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(1.05)
for ci in range(2, 8):
    tbl.columns[ci].width = Inches(1.325)

for ri in range(1, len(tbl.rows)):
    for ci in range(len(tbl.columns)):
        for p in tbl.cell(ri, ci).text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(8)

footnote = next(sh for sh in s10.shapes if sh.has_text_frame and sh.top > Inches(5))
footnote.top = Inches(6.2)
for p in footnote.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(8)

prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
