"""Follow-up: row.height (not shape.height) is what drives PowerPoint COM's
rendered table height for a fresh add_table() table -- shrink each of the 23
rows from 0.27in (deck6/7) to 0.22in (~5.06in total for 23 rows, matching the
8pt body font set in deck7), reposition the table just under the title, and
move the footnote into the freed space above the footer band.
"""

from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
prs = Presentation(PPTX_PATH)

s10 = prs.slides[10]
tbl_shape = next(sh for sh in s10.shapes if sh.has_table)
tbl = tbl_shape.table

tbl_shape.top = Inches(0.75)
tbl_shape.height = Inches(5.06)
for ri in range(len(tbl.rows)):
    tbl.rows[ri].height = Inches(0.22)

footnote = next(sh for sh in s10.shapes if sh.has_text_frame and sh.top > Inches(5))
footnote.top = Inches(5.85)
footnote.height = Inches(0.7)

prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
