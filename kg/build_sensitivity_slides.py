"""
build_sensitivity_slides.py — Add slides 11, 12, 13 to the BCC progress PPTX template.
Creates properly formatted tables for the three sensitivity dimensions.

Run once to set up the template, then populate_sensitivity_slides.py fills the cells.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree


PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
HEADER_BG = '2F5496'
HEADER_FG = 'FFFFFF'


def _set_cell(tbl, row, col, text, bold=False, font_size=Pt(7), color=None):
    cell = tbl.cell(row, col)
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(
            int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _style_header_cell(tbl, row, col):
    cell = tbl.cell(row, col)
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgbClr = etree.SubElement(solidFill,
        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgbClr.set('val', HEADER_BG)


def _add_textbox(slide, text, left, top, width, height, font_size=Pt(10), bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    return txBox


def build():
    prs = Presentation(PPTX_PATH)

    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name and 'blank' in layout.name.lower():
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]

    LEFT = Inches(0.4)
    WIDTH = Inches(11.2)

    # ═══════ Slide 11: Demand Sensitivity ═══════════════════════════════════════
    s11 = prs.slides.add_slide(blank_layout)
    _add_textbox(s11,
        '3  \u2014  Sensitivity: Demand scaling (0.8\u00d7 / 1.0\u00d7 / 1.2\u00d7) \u2014 NO_TSP vs NashGate',
        LEFT, Inches(0.05), Inches(11.5), Inches(0.5), bold=True)
    _add_textbox(s11,
        '[Auto-populated from batch_results.csv]',
        LEFT, Inches(6.3), Inches(11.5), Inches(0.25), font_size=Pt(7))

    nrows, ncols = 22, 10  # 1 header + 21 data rows
    top = Inches(0.55)
    height = Inches(5.65)

    tbl_shape = s11.shapes.add_table(nrows, ncols, LEFT, top, WIDTH, height)
    tbl = tbl_shape.table
    col_widths = [Inches(1.6)] + [Inches(0.95)] * 9
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    # Style header
    for ci in range(ncols):
        _style_header_cell(tbl, 0, ci)

    DEMAND_GROUPS = [('0.8\u00d7', 'NO_TSP_0.8x', 'NashGate_0.8x'),
                     ('1.0\u00d7', 'NO_TSP_1.0x', 'NashGate_1.0x'),
                     ('1.2\u00d7', 'NO_TSP_1.2x', 'NashGate_1.2x')]
    _set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG)
    for gi, (label, _, _) in enumerate(DEMAND_GROUPS):
        bc = 1 + gi * 3
        _set_cell(tbl, 0, bc,     f'NO_TSP {label}', bold=True, color=HEADER_FG)
        _set_cell(tbl, 0, bc + 1, f'NashGate {label}', bold=True, color=HEADER_FG)
        _set_cell(tbl, 0, bc + 2, '\u0394%', bold=True, color=HEADER_FG)

    labels_d = [
        'Bus TT (h)', 'Avg Bus TT (min)', 'Avg Car TT (min)', 'Wtd Avg TT (min)',
        'Total Pass Delay (h)', 'Pax-Equiv Passages', 'Avg Pass Delay (s)',
        'Side Pass Delay (h)', '\u2500\u2500\u2500', 'Vehicle Delay (h)',
        'Total Travel Time (h)', 'Car Dist (km)', 'Bus Dist (km)',
        'Total Vehicles', 'Distinct Buses', 'Flow (veh/h)',
        '\u2500\u2500\u2500', 'Z1 (pax-delay)', 'Z2 (offset)', 'Z4 (TT)', 'Objective',
    ]
    for ri, label in enumerate(labels_d):
        _set_cell(tbl, ri + 1, 0, label, bold=('\u2500' in label), font_size=Pt(7))

    # ═══════ Slide 12: Occupancy Sensitivity ════════════════════════════════════
    s12 = prs.slides.add_slide(blank_layout)
    _add_textbox(s12,
        '3  \u2014  Sensitivity: Occupancy (bus/car pax per veh) \u2014 NashGate',
        LEFT, Inches(0.05), Inches(11.5), Inches(0.5), bold=True)
    _add_textbox(s12,
        '[Auto-populated from batch_results.csv]',
        LEFT, Inches(6.3), Inches(11.5), Inches(0.25), font_size=Pt(7))

    nrows2, ncols2 = 22, 4
    tbl_shape2 = s12.shapes.add_table(nrows2, ncols2, LEFT, top, WIDTH, height)
    tbl2 = tbl_shape2.table
    col_widths2 = [Inches(1.8), Inches(3.1), Inches(3.1), Inches(3.1)]
    for ci, cw in enumerate(col_widths2):
        tbl2.columns[ci].width = cw

    for ci in range(ncols2):
        _style_header_cell(tbl2, 0, ci)

    occ_headers = ['LOW\n(bus 20 / car 1.0)', 'BASE\n(bus 40 / car 1.2)', 'HIGH\n(bus 60 / car 1.5)']
    _set_cell(tbl2, 0, 0, 'KPI', bold=True, color=HEADER_FG)
    for ci, hdr in enumerate(occ_headers):
        _set_cell(tbl2, 0, ci + 1, hdr, bold=True, color=HEADER_FG)

    labels_occ = labels_d  # same
    for ri, label in enumerate(labels_occ):
        _set_cell(tbl2, ri + 1, 0, label, bold=('\u2500' in label), font_size=Pt(7))

    # ═══════ Slide 13: Objective-weight Sensitivity ═════════════════════════════
    s13 = prs.slides.add_slide(blank_layout)
    _add_textbox(s13,
        '3  \u2014  Sensitivity: Objective weights (\u03b1 \u00d7 \u03b2) \u2014 NashGate (\u03b3=0)',
        LEFT, Inches(0.05), Inches(11.5), Inches(0.5), bold=True)
    _add_textbox(s13,
        '[Auto-populated from batch_results.csv]',
        LEFT, Inches(6.3), Inches(11.5), Inches(0.25), font_size=Pt(7))

    nrows3, ncols3 = 13, 10
    tbl_shape3 = s13.shapes.add_table(nrows3, ncols3, LEFT, top, WIDTH, height)
    tbl3 = tbl_shape3.table
    col_widths3 = [Inches(1.5)] + [Inches(1.07)] * 9
    for ci, cw in enumerate(col_widths3):
        tbl3.columns[ci].width = cw

    for ci in range(ncols3):
        _style_header_cell(tbl3, 0, ci)

    wobj_hdrs = ['\u03b1=.7 \u03b2=.1', '\u03b1=.7 \u03b2=.2', '\u03b1=.7 \u03b2=.3',
                 '\u03b1=.8 \u03b2=.1', '\u03b1=.8 \u03b2=.2', '\u03b1=.8 \u03b2=.3',
                 '\u03b1=.9 \u03b2=.1', '\u03b1=.9 \u03b2=.2', '\u03b1=.9 \u03b2=.3']
    _set_cell(tbl3, 0, 0, 'KPI (\u03b3=0)', bold=True, color=HEADER_FG, font_size=Pt(7))
    for ci, hdr in enumerate(wobj_hdrs):
        _set_cell(tbl3, 0, ci + 1, hdr, bold=True, color=HEADER_FG, font_size=Pt(6))

    labels_w = [
        'Total Pass Delay (h)', 'Avg Pass Delay (s)', 'Side Pass Delay (h)',
        'Total Travel Time (h)', 'Total Vehicles', 'Distinct Buses',
        'Flow (veh/h)', 'Z1 (pax-delay)', 'Z2 (offset)',
        'Z3 (lateness)', 'Z4 (TT)', 'Objective',
    ]
    for ri, label in enumerate(labels_w):
        _set_cell(tbl3, ri + 1, 0, label, font_size=Pt(7))

    prs.save(PPTX_PATH)
    print(f'Done. Added slides 11, 12, 13 to {PPTX_PATH}')
    print(f'  Total slides now: {len(prs.slides)}')


if __name__ == '__main__':
    build()
