"""
update_progress_deck13.py — Rebuild slide 8 (Simulation results) table from scratch.

Same structural-rendering bug as slide 9 (fixed in deck12).  The old table was
created by a prior script that left the column XML in a form PowerPoint's renderer
doesn't fully honour, causing the last 2 columns (CellSearch, CellQLearn) to
appear blank in exported renders.
"""

import math
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH  = 'batch_results.csv'

RHO_BUS = 40.0
RHO_CAR = 1.5

HEADER_FILL    = '00467F'
NOPRIORITY_FILL = 'FFF2CC'
BORDER_COLOR   = 'BFBFBF'
BORDER_W       = 6350

# ── Formatting helpers ────────────────────────────────────────────────────────
def f1(x):  return f'{x:,.1f}'
def f2(x):  return f'{x:,.2f}'
def fi(x):  return f'{int(round(x)):,}'

FMT = {
    1: f1, 2: f2, 3: f2, 4: f2, 5: f1, 6: fi,  7: f1,  8: f1,
    10: f1, 11: f1, 12: fi, 13: f2, 14: fi, 15: fi, 16: f1,
    18: f1, 19: f1, 21: f1, 22: f1,
}

# ── Formula computation ───────────────────────────────────────────────────────
def compute_kpi(d):
    r = {}
    r[1]  = d['stats_Net_TotalTT_h_Bus']
    r[2]  = d['stats_Net_TotalTT_h_Bus'] * 60.0 / d['stats_N_DistinctBuses']
    r[3]  = d['stats_Net_TotalTT_h_Car']  * 60.0 / d['stats_N_DistinctCars']
    num4  = (d['stats_Net_TotalTT_h_Car']   * RHO_CAR +
             d['stats_Net_TotalTT_h_Bus']   * RHO_BUS +
             d['stats_Net_TotalTT_h_Truck'] * RHO_CAR) * 60.0
    den4  = (d['stats_N_DistinctCars']   * RHO_CAR +
             d['stats_N_DistinctBuses']  * RHO_BUS +
             d['stats_N_DistinctTrucks'] * RHO_CAR)
    r[4]  = num4 / den4
    r[5]  = d['stats_TotalPassDelay_hrs']
    r[6]  = d['stats_PaxEquivPassages']
    r[7]  = d['stats_AvgPassDelay_s']
    r[8]  = d['stats_SidePassDelay_hrs']
    totveh = (d['stats_N_DistinctCars'] +
              d['stats_N_DistinctBuses'] +
              d['stats_N_DistinctTrucks'])
    r[10] = d['stats_Net_Delay_All'] * totveh / 3600.0
    r[11] = (d['stats_Net_TotalTT_h_Car'] +
             d['stats_Net_TotalTT_h_Bus'] +
             d['stats_Net_TotalTT_h_Truck'])
    r[12] = d['stats_Net_TotalDist_Car']
    r[13] = d['stats_Net_TotalDist_Bus']
    r[14] = totveh
    r[15] = d['stats_N_DistinctBuses']
    r[16] = (d['stats_Net_Flow_Car'] +
             d['stats_Net_Flow_Bus'] +
             d['stats_Net_Flow_Truck'])
    r[18] = d['wobj_Z1_total']       # NaN for NO_TSP
    r[19] = d['wobj_Z2_total']       # 0.0 (bug); NaN for NO_TSP
    r[21] = d['wobj_Z4_total']
    r[22] = d['wobj_objective_total'] # NaN for NO_TSP
    return r

# ── XML helpers (identical to deck12) ────────────────────────────────────────
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def _color_fill(hex_color):
    sf = etree.Element(f'{{{A_NS}}}solidFill')
    etree.SubElement(sf, f'{{{A_NS}}}srgbClr', val=hex_color)
    return sf

def _border_line(tag_name, w=BORDER_W, color=BORDER_COLOR):
    ln = etree.Element(f'{{{A_NS}}}{tag_name}',
                       w=str(w), cap='flat', cmpd='sng', algn='ctr')
    ln.append(_color_fill(color))
    etree.SubElement(ln, f'{{{A_NS}}}prstDash', val='solid')
    etree.SubElement(ln, f'{{{A_NS}}}round')
    etree.SubElement(ln, f'{{{A_NS}}}headEnd', type='none', w='med', len='med')
    etree.SubElement(ln, f'{{{A_NS}}}tailEnd', type='none', w='med', len='med')
    return ln

def apply_cell_style(cell, fill_hex=None, no_fill=False,
                     mar_lr=73152, mar_tb=27432):
    for existing in cell._tc.findall(f'{{{A_NS}}}tcPr'):
        cell._tc.remove(existing)
    tcPr = etree.SubElement(cell._tc, f'{{{A_NS}}}tcPr',
                             marL=str(mar_lr), marR=str(mar_lr),
                             marT=str(mar_tb), marB=str(mar_tb),
                             anchor='ctr')
    for ln_tag in ('lnL', 'lnR', 'lnT', 'lnB'):
        tcPr.append(_border_line(ln_tag))
    if fill_hex:
        tcPr.append(_color_fill(fill_hex))
    elif no_fill:
        etree.SubElement(tcPr, f'{{{A_NS}}}noFill')
    return tcPr

def set_cell_text_simple(cell, text, bold=False, color_hex=None,
                         sz=800, align='ctr', second_para=None):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    pPr = etree.SubElement(p._p, f'{{{A_NS}}}pPr', algn=align)
    p._p.insert(0, pPr)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz / 100)
    run.font.bold = bold
    if color_hex:
        run.font.color.rgb = RGBColor.from_string(color_hex)
    if second_para is not None:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = second_para
        run2.font.size = Pt(sz / 100)
        run2.font.bold = bold
        if color_hex:
            run2.font.color.rgb = RGBColor.from_string(color_hex)

# ── Data ──────────────────────────────────────────────────────────────────────
df = pd.read_csv(CSV_PATH)

CORE_ORDER = [
    'NO_TSP', 'DCTSP_MARL', 'DCTSP_ZIG',
    'DCTSP_BARGAIN_SPM', 'DCTSP_MP_ECTM', 'DCTSP_BXT',
]
COL_LABELS = ['NoPriority', 'CPD-QL', 'WaveGate', 'NashGate', 'CellSearch', 'CellQLearn']

d_core   = (df[df['run_experiment'].isin(CORE_ORDER)]
              .set_index('run_experiment').loc[CORE_ORDER])
kpi_core = compute_kpi(d_core)

# NoPriority column offset index and row indices where it shows "—"
NOPRIORITY_CI    = 0                 # col_offset 0 = NO_TSP
NOPRIORITY_SKIP  = {18, 19, 20, 21} # rows (1-based) where NoPriority shows "—"

# ── Row definitions ───────────────────────────────────────────────────────────
ROWS_DATA = [
    ('Total bus travel time',                    '[bus-h]',         1),
    ('Average bus travel time',                  '[min/bus]',       2),
    ('Total car travel time',                    '[min/car]',       3),
    ('Total passenger travel time',              '[min/pax]',       4),
    ('Total passenger delay',                    '[pax-h]',         5),
    ('Total serviced passengers',                '[pax]',           6),
    ('Average passenger delay',                  '[sec/pax]',       7),
    ('Side-street total passenger delay',        '[pax-h]',         8),
    ('Side-street average passenger delay',      '[sec/pax]',       'dash'),
    ('Total system delay',                       '[veh-h]',         10),
    ('Total hours of travel',                    '[veh-h]',         11),
    ('Total VKT – car',                     '[veh-km]',        12),
    ('Total VKT – bus',                     '[veh-km]',        13),
    ('Total vehicles in system',                 '[veh (count)]',   14),
    ('Total vehicles in system (Bus)',           '[veh (count)]',   15),
    ('Throughput – main',                   '[veh/h]',         16),
    ('Throughput – side',                   '[veh/h]',         'dash'),
    ('Weighted pax-delay objective (Z1)',        '[a.u.]',          18),
    ('Offset-correction objective (Z2)',         '[a.u.]',          'z2star'),
    ('Raw travel-time objective (Z4)',           '[a.u.]',          21),
    ('Total weighted objective',                 '[a.u.]',          22),
]

# ── Rebuild slide 8 table ─────────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
s8  = prs.slides[8]

old_sh = next(sh for sh in s8.shapes if sh.has_table)
old_left   = old_sh.left
old_top    = old_sh.top
old_width  = old_sh.width
old_height = old_sh.height
old_sh._element.getparent().remove(old_sh._element)

N_ROWS = 1 + len(ROWS_DATA)   # 22
N_COLS = 2 + len(CORE_ORDER)  # 8

tsh = s8.shapes.add_table(N_ROWS, N_COLS, old_left, old_top, old_width, old_height)
tbl = tsh.table

# Column widths (preserved from old table)
tbl.columns[0].width = Inches(2.7)
tbl.columns[1].width = Inches(0.85)
for ci in range(2, N_COLS):
    tbl.columns[ci].width = Inches(1.342)

# Row heights (uniform)
row_h = old_height // N_ROWS
for ri in range(N_ROWS):
    tbl.rows[ri].height = row_h

# ── Header row ────────────────────────────────────────────────────────────────
apply_cell_style(tbl.cell(0, 0), fill_hex=HEADER_FILL)
set_cell_text_simple(tbl.cell(0, 0), 'KPI', bold=True, color_hex='FFFFFF', align='l')

apply_cell_style(tbl.cell(0, 1), fill_hex=HEADER_FILL)
set_cell_text_simple(tbl.cell(0, 1), 'unit', bold=True, color_hex='FFFFFF')

for ci, label in enumerate(COL_LABELS):
    col = 2 + ci
    apply_cell_style(tbl.cell(0, col), fill_hex=HEADER_FILL)
    set_cell_text_simple(tbl.cell(0, col), label, bold=True, color_hex='FFFFFF')

# ── Data rows ─────────────────────────────────────────────────────────────────
for ri, (label, unit, fkey) in enumerate(ROWS_DATA):
    trow = ri + 1

    apply_cell_style(tbl.cell(trow, 0), no_fill=True)
    set_cell_text_simple(tbl.cell(trow, 0), label, align='l')

    apply_cell_style(tbl.cell(trow, 1), no_fill=True)
    set_cell_text_simple(tbl.cell(trow, 1), unit)

    for ci in range(len(CORE_ORDER)):
        col = 2 + ci
        is_nopriority = (ci == NOPRIORITY_CI)
        # NoPriority column uses yellow fill
        cell_fill = NOPRIORITY_FILL if is_nopriority else None
        apply_cell_style(tbl.cell(trow, col),
                         fill_hex=cell_fill,
                         no_fill=(not is_nopriority))
        # Determine text
        if fkey == 'dash':
            text = '—'
        elif fkey == 'z2star':
            text = '—' if is_nopriority else '0.0*'
        elif is_nopriority and trow in NOPRIORITY_SKIP:
            text = '—'
        else:
            val = kpi_core[fkey].iloc[ci]
            if math.isnan(val):
                text = '—'
            else:
                text = FMT[fkey](val)
        set_cell_text_simple(tbl.cell(trow, col), text)

prs.save(PPTX_PATH)
print('Saved', PPTX_PATH)
