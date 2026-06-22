"""
rebuild_sensitivity_slides.py — Find & rebuild the 3 sensitivity slides in the PPTX.
Each slide gets a dedicated NO_TSP comparison column.

Slide layouts:
  Demand (22r x 10c):    KPI | NO_TSP 0.8x | NashGate 0.8x | Δ% | NO_TSP 1.0x | ... | NO_TSP 1.2x | NashGate 1.2x | Δ%
  Occupancy (22r x 7c):  KPI | NO_TSP | LOW | Δ% | BASE | Δ% | HIGH | Δ%
  Wobj (13r x 11c):      KPI | NO_TSP | 9 weight combos
"""
import shutil, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
HEADER_BG = '2F5496'
HEADER_FG = 'FFFFFF'
LEFT, TOP, WIDTH, HEIGHT = Inches(0.4), Inches(0.55), Inches(11.2), Inches(4.0)  # shorter: room for chart


def set_cell(tbl, r, c, text, bold=False, sz=Pt(7), color=None):
    cell = tbl.cell(r, c)
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = sz
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(int(color[0:2],16), int(color[2:4],16), int(color[4:6],16))

def hdr_fill(tbl):
    for ci in range(len(tbl.columns)):
        tc = tbl.cell(0, ci)._tc.get_or_add_tcPr()
        sf = etree.SubElement(tc, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr').set('val', HEADER_BG)

def remove_table(slide):
    for sh in list(slide.shapes):
        if sh.has_table:
            sh._element.getparent().remove(sh._element)

def add_table(slide, nr, nc, col_widths):
    ts = slide.shapes.add_table(nr, nc, LEFT, TOP, WIDTH, HEIGHT)
    tbl = ts.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    return tbl

def slide_text(slide):
    return ' '.join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)

def _move_footnote_down(slide):
    """Move footnote text boxes down to make room for chart (at 4.65-6.65 inch)."""
    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table and sh.top > Inches(4.0):
            sh.top = Inches(6.75)
            sh.height = Inches(0.35)

def rebuild():
    prs = Presentation(PPTX_PATH)

    # Find slides by title content
    d_idx = o_idx = w_idx = None
    for i, s in enumerate(prs.slides):
        txt = slide_text(s)
        if 'Demand scaling' in txt:
            d_idx = i
        elif 'Occupancy' in txt and 'bus' in txt.lower():
            o_idx = i
        elif 'Objective weights' in txt:
            w_idx = i
    print(f'Found slides: demand={d_idx} occ={o_idx} wobj={w_idx}')

    # Common data row labels
    LABELS_22 = [
        'Bus TT (h)', 'Avg Bus TT (min)', 'Avg Car TT (min)', 'Wtd Avg TT (min)',
        'Total Pass Delay (h)', 'Pax-Equiv Passages', 'Avg Pass Delay (s)',
        'Side Pass Delay (h)', '\u2500\u2500\u2500', 'Vehicle Delay (h)',
        'Total Travel Time (h)', 'Car Dist (km)', 'Bus Dist (km)',
        'Total Vehicles', 'Distinct Buses', 'Flow (veh/h)',
        '\u2500\u2500\u2500', 'Z1: wtd pax delay (pax\u00b7s)', 'Z2: offset-corr (s)', 'Z4: corridor TT (veh\u00b7h)', 'Obj = \u03b1Z1+\u03b2Z2+\u03b3Z3',
    ]
    LABELS_WOBJ = [
        'Total Pass Delay (h)', 'Avg Pass Delay (s)', 'Side Pass Delay (h)',
        'Total Travel Time (h)', 'Total Vehicles', 'Distinct Buses', 'Flow (veh/h)',
        'Z1: wtd pax delay (pax\u00b7s)', 'Z2: offset-corr (s)', 'Z3: bus lateness \u03c3\u207a (s)',
        'Z4: corridor TT (veh\u00b7h)', 'Obj = \u03b1Z1+\u03b2Z2+\u03b3Z3',
    ]

    # ── DEMAND SLIDE ──
    if d_idx is not None:
        s = prs.slides[d_idx]
        remove_table(s)
        tbl = add_table(s, 22, 10, [1.6] + [0.95]*9)
        hdr_fill(tbl)
        set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG)
        for gi, lbl in enumerate(['0.8\u00d7', '1.0\u00d7', '1.2\u00d7']):
            bc = 1 + gi * 3
            set_cell(tbl, 0, bc,     f'NO_TSP {lbl}', bold=True, color=HEADER_FG)
            set_cell(tbl, 0, bc + 1, f'NashGate {lbl}', bold=True, color=HEADER_FG)
            set_cell(tbl, 0, bc + 2, '\u0394%', bold=True, color=HEADER_FG)
        for ri, lbl in enumerate(LABELS_22):
            set_cell(tbl, ri + 1, 0, lbl, bold=('\u2500' in lbl), sz=Pt(7))
        _move_footnote_down(s)
        print('  Demand slide: 22r x 10c OK')

    # ── OCCUPANCY SLIDE ──
    if o_idx is not None:
        s = prs.slides[o_idx]
        remove_table(s)
        tbl = add_table(s, 22, 8, [1.5, 1.3, 1.4, 0.9, 1.4, 0.9, 1.4, 0.9])
        hdr_fill(tbl)
        set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, 1, 'NO_TSP', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, 2, 'LOW\n(20/1.0)', bold=True, color=HEADER_FG, sz=Pt(7))
        set_cell(tbl, 0, 3, '\u0394%', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, 4, 'BASE\n(40/1.2)', bold=True, color=HEADER_FG, sz=Pt(7))
        set_cell(tbl, 0, 5, '\u0394%', bold=True, color=HEADER_FG)
        set_cell(tbl, 0, 6, 'HIGH\n(60/1.5)', bold=True, color=HEADER_FG, sz=Pt(7))
        set_cell(tbl, 0, 7, '\u0394%', bold=True, color=HEADER_FG)
        for ri, lbl in enumerate(LABELS_22):
            set_cell(tbl, ri + 1, 0, lbl, bold=('\u2500' in lbl), sz=Pt(7))
        _move_footnote_down(s)
        print('  Occupancy slide: 22r x 8c OK')

    # ── OBJ-WEIGHT SLIDE (13r x 7c: KPI + NO_TSP + 5 combos) ──
    if w_idx is not None:
        s = prs.slides[w_idx]
        remove_table(s)
        tbl = add_table(s, 13, 7, [1.5, 1.3, 1.6, 1.6, 1.6, 1.6, 1.6])
        hdr_fill(tbl)
        set_cell(tbl, 0, 0, 'KPI', bold=True, color=HEADER_FG, sz=Pt(7))
        set_cell(tbl, 0, 1, 'NO_TSP', bold=True, color=HEADER_FG, sz=Pt(7))
        whdrs = ['EQ Z1+Z2\n\u03b1=.5 \u03b2=.5', 'EQ ALL3\n\u03b1=.33 \u03b2=.33 \u03b3=.33',
                 'Only Z1\n\u03b1=1 \u03b2=0 \u03b3=0', 'Only Z2\n\u03b1=0 \u03b2=1 \u03b3=0',
                 'Only TT(Z4)\n\u03b4=1']
        for ci, hdr in enumerate(whdrs):
            set_cell(tbl, 0, ci + 2, hdr, bold=True, color=HEADER_FG, sz=Pt(6))
        for ri, lbl in enumerate(LABELS_WOBJ):
            set_cell(tbl, ri + 1, 0, lbl, sz=Pt(7))
        _move_footnote_down(s)
        print('  Wobj slide: 13r x 7c OK')

    prs.save(PPTX_PATH)
    print(f'Saved {PPTX_PATH} ({len(prs.slides)} slides total)')
    shutil.copy(PPTX_PATH, PPTX_PATH.replace('_todo.pptx', '_fake.pptx'))
    print(f'Copied to fake PPTX')


if __name__ == '__main__':
    rebuild()
