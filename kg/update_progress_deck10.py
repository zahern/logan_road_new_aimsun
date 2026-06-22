"""Follow-up #2: deck9's row.height=0.22in still rendered ~0.235in/row
(23 x 0.235 ~= 5.4in), because PowerPoint's content-driven minimum row height
(8pt line height + default 0.05in top/bottom cell margins) exceeds 0.22in,
so the table still overran the footnote at top=5.85in.

Shrink the per-cell top/bottom margins to 0.02in (from the 0.05in default),
which lowers the content-driven floor to ~0.17in, and request row
height=0.18in. With 23 rows this should land close to ~4.1-4.4in total,
comfortably clear of the footer band. Move the footnote down into the
freed space and restore its font to 9pt.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
prs = Presentation(PPTX_PATH)

s10 = prs.slides[10]
tbl_shape = next(sh for sh in s10.shapes if sh.has_table)
tbl = tbl_shape.table

for ri in range(len(tbl.rows)):
    tbl.rows[ri].height = Inches(0.18)
    for ci in range(len(tbl.columns)):
        cell = tbl.cell(ri, ci)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)

tbl_shape.top = Inches(0.75)
tbl_shape.height = Inches(0.18 * 23)

footnote = next(sh for sh in s10.shapes if sh.has_text_frame and sh.top > Inches(4))
footnote.top = Inches(5.0)
footnote.height = Inches(1.8)
for p in footnote.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(9)

prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
