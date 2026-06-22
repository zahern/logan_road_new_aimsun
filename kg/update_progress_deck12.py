"""
update_progress_deck12.py — Rebuild slide 9 (WOBJ sweep) table from scratch.

The existing table has a malformed XML structure (duplicate a16:colId values in
tblGrid from an earlier session's script) that causes PowerPoint to render only
the first 6 of 11 columns.  Replacing the shape with a fresh add_table() call
produces well-formed XML identical in structure to slide 10 (Preview), which
renders correctly.
"""

import math
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH  = 'batch_results.csv'

RHO_BUS = 40.0
RHO_CAR = 1.5

HEADER_FILL  = '00467F'
BORDER_COLOR = 'BFBFBF'
BORDER_W     = 6350   # 0.5pt in EMU (hundredths of a point × 127)

# ── Formatting helpers ────────────────────────────────────────────────────────
def f1(x):  return f'{x:,.1f}'
def f2(x):  return f'{x:,.2f}'
def fi(x):  return f'{int(round(x)):,}'

FMT = {
    1: f1, 2: f2, 3: f2, 4: f2, 5: f1, 6: fi,  7: f1,  8: f1,
    10: f1, 11: f1, 12: fi, 13: f2, 14: fi, 15: fi, 16: f1,
    18: f1, 19: f1, 21: f1, 22: f1,
}

# ── Formula computation ───────────────────────────────────────────────────────
def compute_kpi(d):
    r = {}
    r[1]  = d['stats_Net_TotalTT_h_Bus']
    r[2]  = d['stats_Net_TotalTT_h_Bus'] * 60.0 / d['stats_N_DistinctBuses']
    r[3]  = d['stats_Net_TotalTT_h_Car']  * 60.0 / d['stats_N_DistinctCars']
    num4  = (d['stats_Net_TotalTT_h_Car']   * RHO_CAR +
             d['stats_Net_TotalTT_h_Bus']   * RHO_BUS +
             d['stats_Net_TotalTT_h_Truck'] * RHO_CAR) * 60.0
    den4  = (d['stats_N_DistinctCars']   * RHO_CAR +
             d['stats_N_DistinctBuses']  * RHO_BUS +
             d['stats_N_DistinctTrucks'] * RHO_CAR)
    r[4]  = num4 / den4
    r[5]  = d['stats_TotalPassDelay_hrs']
    r[6]  = d['stats_PaxEquivPassages']
    r[7]  = d['stats_AvgPassDelay_s']
    r[8]  = d['stats_SidePassDelay_hrs']
    totveh = (d['stats_N_DistinctCars'] +
              d['stats_N_DistinctBuses'] +
              d['stats_N_DistinctTrucks'])
    r[10] = d['stats_Net_Delay_All'] * totveh / 3600.0
    r[11] = (d['stats_Net_TotalTT_h_Car'] +
             d['stats_Net_TotalTT_h_Bus'] +
             d['stats_Net_TotalTT_h_Truck'])
    r[12] = d['stats_Net_TotalDist_Car']
    r[13] = d['stats_Net_TotalDist_Bus']
    r[14] = totveh
    r[15] = d['stats_N_DistinctBuses']
    r[16] = (d['stats_Net_Flow_Car'] +
             d['stats_Net_Flow_Bus'] +
             d['stats_Net_Flow_Truck'])
    r[18] = d['wobj_Z1_total']
    r[19] = d['wobj_Z2_total']   # = 0.0 (bug); will show '0.0*'
    r[21] = d['wobj_Z4_total']
    r[22] = d['wobj_objective_total']
    return r

# ── XML helpers ───────────────────────────────────────────────────────────────
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def _el(tag, attrib=None, nsmap=None):
    return etree.SubElement if False else etree.Element(
        f'{{{A_NS}}}{tag}', attrib=attrib or {})

def _sub(parent, tag, attrib=None):
    el = etree.SubElement(parent, f'{{{A_NS}}}{tag}', attrib=attrib or {})
    return el

def _color_fill(hex_color):
    """Return an <a:solidFill> element."""
    sf = etree.Element(f'{{{A_NS}}}solidFill')
    clr = etree.SubElement(sf, f'{{{A_NS}}}srgbClr', val=hex_color)
    return sf

def _border_line(parent_tag_name, w=BORDER_W, color=BORDER_COLOR):
    """Return an <a:lnX> border element."""
    ln = etree.Element(f'{{{A_NS}}}{parent_tag_name}',
                       w=str(w), cap='flat', cmpd='sng', algn='ctr')
    ln.append(_color_fill(color))
    pd_ = etree.SubElement(ln, f'{{{A_NS}}}prstDash', val='solid')
    etree.SubElement(ln, f'{{{A_NS}}}round')
    etree.SubElement(ln, f'{{{A_NS}}}headEnd', type='none', w='med', len='med')
    etree.SubElement(ln, f'{{{A_NS}}}tailEnd', type='none', w='med', len='med')
    return ln

def apply_cell_style(cell, fill_hex=None, no_fill=False,
                     mar_lr=73152, mar_tb=27432):
    """Add borders and optional fill to a table cell via direct XML."""
    # Remove any existing tcPr (add_table() inserts an empty one by default)
    for existing in cell._tc.findall(f'{{{A_NS}}}tcPr'):
        cell._tc.remove(existing)
    tcPr = etree.SubElement(cell._tc, f'{{{A_NS}}}tcPr',
                             marL=str(mar_lr), marR=str(mar_lr),
                             marT=str(mar_tb), marB=str(mar_tb),
                             anchor='ctr')
    for ln_tag in ('lnL', 'lnR', 'lnT', 'lnB'):
        tcPr.append(_border_line(ln_tag))
    if fill_hex:
        tcPr.append(_color_fill(fill_hex))
    elif no_fill:
        etree.SubElement(tcPr, f'{{{A_NS}}}noFill')
    return tcPr

def set_cell_text_simple(cell, text, bold=False, color_hex=None,
                         sz=800, align='ctr', second_para=None):
    """Populate a cell's text frame (single paragraph, one run)."""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # Set paragraph alignment
    pPr = etree.SubElement(p._p, f'{{{A_NS}}}pPr', algn=align)
    p._p.insert(0, pPr)  # pPr must be first child of p
    # Build run
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz / 100)   # sz=800 → 8pt
    run.font.bold = bold
    if color_hex:
        run.font.color.rgb = RGBColor.from_string(color_hex)
    if second_para is not None:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = second_para
        run2.font.size = Pt(sz / 100)
        run2.font.bold = bold
        if color_hex:
            run2.font.color.rgb = RGBColor.from_string(color_hex)

# ── Data ──────────────────────────────────────────────────────────────────────
df = pd.read_csv(CSV_PATH)

WOBJ_ORDER = [
    'WOBJ_SWEEP_A07_B01', 'WOBJ_SWEEP_A07_B02', 'WOBJ_SWEEP_A07_B03',
    'WOBJ_SWEEP_A08_B01', 'WOBJ_SWEEP_A08_B02', 'WOBJ_SWEEP_A08_B03',
    'WOBJ_SWEEP_A09_B01', 'WOBJ_SWEEP_A09_B02', 'WOBJ_SWEEP_A09_B03',
]
d_wobj   = (df[df['run_experiment'].isin(WOBJ_ORDER)]
              .set_index('run_experiment').loc[WOBJ_ORDER])
kpi_wobj = compute_kpi(d_wobj)

# ── Row definitions ───────────────────────────────────────────────────────────
# KPI label rows (label, unit, formula_key or 'dash' or 'z2star')
ROWS_DATA = [
    ('Total bus travel time',                    '[bus-h]',         1),
    ('Average bus travel time',                  '[min/bus]',       2),
    ('Total car travel time',                    '[min/car]',       3),
    ('Total passenger travel time',              '[min/pax]',       4),
    ('Total passenger delay',                    '[pax-h]',         5),
    ('Total serviced passengers',                '[pax]',           6),
    ('Average passenger delay',                  '[sec/pax]',       7),
    ('Side-street total passenger delay',        '[pax-h]',         8),
    ('Side-street average passenger delay',      '[sec/pax]',       'dash'),
    ('Total system delay',                       '[veh-h]',         10),
    ('Total hours of travel',                    '[veh-h]',         11),
    ('Total VKT – car',                     '[veh-km]',        12),
    ('Total VKT – bus',                     '[veh-km]',        13),
    ('Total vehicles in system',                 '[veh (count)]',   14),
    ('Total vehicles in system (Bus)',           '[veh (count)]',   15),
    ('Throughput – main',                   '[veh/h]',         16),
    ('Throughput – side',                   '[veh/h]',         'dash'),
    ('Weighted pax-delay objective (Z1)',        '[a.u.]',          18),
    ('Offset-correction objective (Z2)',         '[a.u.]',          'z2star'),
    ('Raw travel-time objective (Z4)',           '[a.u.]',          21),
    ('Total weighted objective',                 '[a.u.]',          22),
]

# 9 column headers (α=X.Y label, β=Z.W label as separate paragraphs)
COL_HEADERS = [
    ('α=0.7', 'β=0.1'), ('α=0.7', 'β=0.2'), ('α=0.7', 'β=0.3'),
    ('α=0.8', 'β=0.1'), ('α=0.8', 'β=0.2'), ('α=0.8', 'β=0.3'),
    ('α=0.9', 'β=0.1'), ('α=0.9', 'β=0.2'), ('α=0.9', 'β=0.3'),
]

# ── Rebuild slide 9 table ─────────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
s9  = prs.slides[9]

# Remove old (broken) table shape
old_sh = next(sh for sh in s9.shapes if sh.has_table)
old_left   = old_sh.left
old_top    = old_sh.top
old_width  = old_sh.width
old_height = old_sh.height
old_sh._element.getparent().remove(old_sh._element)

N_ROWS = 1 + len(ROWS_DATA)  # 22
N_COLS = 2 + len(WOBJ_ORDER) # 11

tsh = s9.shapes.add_table(N_ROWS, N_COLS, old_left, old_top, old_width, old_height)
tbl = tsh.table

# Column widths
tbl.columns[0].width = Inches(2.4)
tbl.columns[1].width = Inches(0.55)
for ci in range(2, N_COLS):
    tbl.columns[ci].width = Inches(0.961)

# Row heights (uniform)
row_h = old_height // N_ROWS
for ri in range(N_ROWS):
    tbl.rows[ri].height = row_h

# ── Header row (row 0) ────────────────────────────────────────────────────────
# KPI
apply_cell_style(tbl.cell(0, 0), fill_hex=HEADER_FILL)
set_cell_text_simple(tbl.cell(0, 0), 'KPI', bold=True, color_hex='FFFFFF', align='l')

# unit
apply_cell_style(tbl.cell(0, 1), fill_hex=HEADER_FILL)
set_cell_text_simple(tbl.cell(0, 1), 'unit', bold=True, color_hex='FFFFFF')

# 9 data columns
for ci, (line1, line2) in enumerate(COL_HEADERS):
    col = 2 + ci
    apply_cell_style(tbl.cell(0, col), fill_hex=HEADER_FILL)
    set_cell_text_simple(tbl.cell(0, col), line1, bold=True,
                         color_hex='FFFFFF', second_para=line2)

# ── Data rows (rows 1-21) ─────────────────────────────────────────────────────
for ri, (label, unit, fkey) in enumerate(ROWS_DATA):
    trow = ri + 1

    # KPI label cell (left-aligned)
    apply_cell_style(tbl.cell(trow, 0), no_fill=True)
    set_cell_text_simple(tbl.cell(trow, 0), label, align='l')

    # unit cell
    apply_cell_style(tbl.cell(trow, 1), no_fill=True)
    set_cell_text_simple(tbl.cell(trow, 1), unit)

    # 9 data cells
    for ci in range(len(WOBJ_ORDER)):
        col = 2 + ci
        apply_cell_style(tbl.cell(trow, col), no_fill=True)
        if fkey == 'dash':
            text = '—'
        elif fkey == 'z2star':
            text = '0.0*'
        else:
            val = kpi_wobj[fkey].iloc[ci]
            text = FMT[fkey](val) if not math.isnan(val) else '—'
        set_cell_text_simple(tbl.cell(trow, col), text)

prs.save(PPTX_PATH)
print('Saved', PPTX_PATH)
