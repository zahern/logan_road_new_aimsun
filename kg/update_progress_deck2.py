"""Follow-up one-off script (run after update_progress_deck.py):

  - Slide "2 - ... Simulation results" (the KPI table, 6 strategy columns):
    add 4 new rows for the weighted-objective components (Z1/Z2/Z4/total),
    plus a "N of 58 configurations" note in the title.
  - New slide "2 - ... Sensitivity of objectives (WOBJ_ALPHA x WOBJ_BETA)":
    same table (20 rows x KPI/unit + 9 WOBJ_SWEEP columns), inserted right
    after the KPI-table slide.
  - Slide "1.2 - Sensitivity sweeps configured this period": add a caption
    summarising the total configuration counts.

Source data: batch_results.csv.
"""

from copy import deepcopy

import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH = 'batch_results.csv'

# =============================================================================
# 1. Load data
# =============================================================================
df = pd.read_csv(CSV_PATH)
RHO_BUS, RHO_CAR = 40.0, 1.5

STRATS = [
    ("NoPriority", "NO_TSP"),
    ("CPD-QL",     "DCTSP_MARL"),
    ("WaveGate",   "DCTSP_ZIG"),
    ("NashGate",   "DCTSP_BARGAIN_SPM"),
    ("CellSearch", "DCTSP_MP_ECTM"),
    ("CellQLearn", "DCTSP_BXT"),
]
rdata = {}
for label, exp in STRATS:
    rdata[label] = df[df['run_experiment'] == exp].iloc[0]

WOBJ_ALPHAS = [0.7, 0.8, 0.9]
WOBJ_BETAS = [0.1, 0.2, 0.3]
wobj_rdata = {}
for a in WOBJ_ALPHAS:
    for b in WOBJ_BETAS:
        name = f"WOBJ_SWEEP_A{a:.1f}_B{b:.1f}".replace(".", "")
        wobj_rdata[(a, b)] = df[df['run_experiment'] == name].iloc[0]


def fmt(v, nd=1):
    if pd.isna(v):
        return "—"
    return f"{v:,.{nd}f}"


def base_kpi_rows(rows_dict):
    """The same 16 KPI rows used on the original 'Simulation results' slide,
    computed for an arbitrary dict of {label: row}."""
    out = []

    out.append(("Total bus travel time", "[min/bus]",
                 [fmt(r["stats_Net_TotalTT_h_Bus"] * 60.0 / r["stats_N_DistinctBuses"], 2) for r in rows_dict.values()]))

    out.append(("Total car travel time", "[min/car]",
                 [fmt(r["stats_Net_TotalTT_h_Car"] * 60.0 / r["stats_N_DistinctCars"], 2) for r in rows_dict.values()]))

    def pax_tt(r):
        num = (r["stats_Net_TotalTT_h_Car"] * RHO_CAR
               + r["stats_Net_TotalTT_h_Bus"] * RHO_BUS
               + r["stats_Net_TotalTT_h_Truck"] * RHO_CAR) * 60.0
        den = (r["stats_N_DistinctCars"] * RHO_CAR
               + r["stats_N_DistinctBuses"] * RHO_BUS
               + r["stats_N_DistinctTrucks"] * RHO_CAR)
        return num / den

    out.append(("Total passenger travel time", "[min/pax]",
                 [fmt(pax_tt(r), 2) for r in rows_dict.values()]))

    out.append(("Total passenger delay", "[pax-h]",
                 [fmt(r["stats_TotalPassDelay_hrs"], 1) for r in rows_dict.values()]))

    out.append(("Total serviced passengers", "[pax]",
                 [fmt(r["stats_PaxEquivPassages"], 0) for r in rows_dict.values()]))

    out.append(("Average passenger delay", "[sec/pax]",
                 [fmt(r["stats_AvgPassDelay_s"], 1) for r in rows_dict.values()]))

    out.append(("Side-street total passenger delay", "[pax-h]",
                 [fmt(r["stats_SidePassDelay_hrs"], 1) for r in rows_dict.values()]))

    out.append(("Side-street average passenger delay", "[sec/pax]", ["—"] * len(rows_dict)))

    def total_veh(r):
        return r["stats_N_DistinctCars"] + r["stats_N_DistinctBuses"] + r["stats_N_DistinctTrucks"]

    def system_delay_veh_h(r):
        return r["stats_Net_Delay_All"] * total_veh(r) / 3600.0

    out.append(("Total system delay", "[veh-h]",
                 [fmt(system_delay_veh_h(r), 1) for r in rows_dict.values()]))

    out.append(("Total hours of travel", "[veh-h]",
                 [fmt(r["stats_Net_TotalTT_h_Car"] + r["stats_Net_TotalTT_h_Bus"] + r["stats_Net_TotalTT_h_Truck"], 1)
                  for r in rows_dict.values()]))

    out.append(("Total VKT – car", "[veh-km]",
                 [fmt(r["stats_Net_TotalDist_Car"], 0) for r in rows_dict.values()]))

    out.append(("Total VKT – bus", "[veh-km]",
                 [fmt(r["stats_Net_TotalDist_Bus"], 2) for r in rows_dict.values()]))

    out.append(("Total vehicles in system", "[veh]",
                 [fmt(total_veh(r), 0) for r in rows_dict.values()]))

    out.append(("Total vehicles in system (Bus)", "[veh]",
                 [fmt(r["stats_N_DistinctBuses"], 0) for r in rows_dict.values()]))

    out.append(("Throughput – main", "[veh/h]",
                 [fmt(r["stats_Net_Flow_Car"] + r["stats_Net_Flow_Bus"] + r["stats_Net_Flow_Truck"], 1)
                  for r in rows_dict.values()]))

    out.append(("Throughput – side", "[veh/h]", ["—"] * len(rows_dict)))

    assert len(out) == 16
    return out


def objective_rows(rows_dict):
    """4 rows for the weighted-objective components Z1/Z2/Z4/total."""
    return [
        ("Weighted pax-delay objective (Z1)", "[a.u.]",
         [fmt(r.get("wobj_Z1_total"), 1) for r in rows_dict.values()]),
        ("Offset-correction objective (Z2)", "[a.u.]",
         [fmt(r.get("wobj_Z2_total"), 1) for r in rows_dict.values()]),
        ("Raw travel-time objective (Z4)", "[a.u.]",
         [fmt(r.get("wobj_Z4_total"), 1) for r in rows_dict.values()]),
        ("Total weighted objective", "[a.u.]",
         [fmt(r.get("wobj_objective_total"), 1) for r in rows_dict.values()]),
    ]


full_rows_strats = base_kpi_rows(rdata) + objective_rows(rdata)        # 20 rows, 6 strategy cols
full_rows_wobj = base_kpi_rows(wobj_rdata) + objective_rows(wobj_rdata)  # 20 rows, 9 WOBJ-config cols
assert len(full_rows_strats) == 20
assert len(full_rows_wobj) == 20

# =============================================================================
# 2. Open presentation
# =============================================================================
prs = Presentation(PPTX_PATH)


def _set_cell_lines(cell, lines, size_pt=9, bold=False):
    """Set a table cell's text to one paragraph per entry in `lines`,
    applying explicit font size/bold to every run."""
    tf = cell.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    if not p0.runs:
        p0.add_run()
    p0.runs[0].text = lines[0]
    p0.runs[0].font.size = Pt(size_pt)
    p0.runs[0].font.bold = bold
    for line in lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.bold = bold


def _set_cell(cell, text, size_pt=9, bold=False):
    _set_cell_lines(cell, [text], size_pt=size_pt, bold=bold)


# Slide 6 (0-indexed) = "2 - ... Simulation results - Kelvin Grove Rd ... KPI table"
s6 = prs.slides[6]
shapes6 = {sh.name: sh for sh in s6.shapes}
tbl6_shape = shapes6["table"]
tbl6 = tbl6_shape.table
_tbl6 = tbl6._tbl

N_COLS = 8  # KPI, unit, 6 strategies

# ── Expand slide 6's table from 17 -> 21 rows (add the 4 objective rows) ───
last_tr = _tbl6.tr_lst[-1]
for _ in range(4):
    _tbl6.append(deepcopy(last_tr))

tbl6 = tbl6_shape.table  # re-fetch after XML mutation
N_ROWS = len(tbl6.rows)  # 21
assert N_ROWS == 21

# Even row heights across the (now taller) table.
for row in tbl6.rows:
    row.height = Inches(6.21 / N_ROWS)

for ri, (label, unit, values) in enumerate(objective_rows(rdata), start=17):
    _set_cell(tbl6.cell(ri, 0), label, size_pt=9, bold=False)
    _set_cell(tbl6.cell(ri, 1), unit, size_pt=9, bold=False)
    for ci, val in enumerate(values, start=2):
        _set_cell(tbl6.cell(ri, ci), val, size_pt=9, bold=False)

# Note the configuration count in the title.
s6_title = shapes6["title"]
title_runs6 = s6_title.text_frame.paragraphs[0].runs
title_runs6[2].text += "  —  6 of 58 total configurations run to date"

# =============================================================================
# 3. New slide "2 - ... Sensitivity of objectives (WOBJ_ALPHA x WOBJ_BETA)"
# =============================================================================
new_slide = prs.slides.add_slide(prs.slide_layouts[8])  # "3_Standard slide"

from pptx.enum.shapes import PP_PLACEHOLDER
for shape in list(new_slide.shapes):
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
        shape._element.getparent().remove(shape._element)

# -- Title: clone slide 6's title shape and edit the topic run. --
new_title_el = deepcopy(s6_title._element)
new_slide.shapes._spTree.append(new_title_el)
new_title = new_slide.shapes[-1]
cNvPr = new_title._element.find('.//' + qn('p:cNvPr'))
cNvPr.set('id', '101')
new_title_runs = new_title.text_frame.paragraphs[0].runs
new_title_runs[2].text = (
    "Sensitivity of objectives — WOBJ_ALPHA × WOBJ_BETA sweep "
    "(GLOBAL_REWARD + KALMAN coordination)  —  9 of 58 total configurations run to date"
)

# -- Table: clone slide 6's (now 21-row x 8-col) table, expand to 11 cols. --
new_tbl_el = deepcopy(tbl6_shape._element)
new_slide.shapes._spTree.append(new_tbl_el)
new_tbl_shape = new_slide.shapes[-1]
cNvPr2 = new_tbl_shape._element.find('.//' + qn('p:cNvPr'))
cNvPr2.set('id', '102')

new_tbl = new_tbl_shape.table
_new_tbl = new_tbl._tbl
tblGrid = _new_tbl.tblGrid

last_gridCol = tblGrid.gridCol_lst[-1]
for _ in range(3):
    tblGrid.append(deepcopy(last_gridCol))

for tr in _new_tbl.tr_lst:
    tcs = tr.tc_lst
    last_tc = tcs[-1]
    for _ in range(3):
        tr.append(deepcopy(last_tc))

new_tbl = new_tbl_shape.table  # re-fetch
N_COLS2 = len(new_tbl.columns)
assert N_COLS2 == 11

# Column widths: KPI | unit | 9 WOBJ-config columns, total ~11.6".
col_widths = [2.4, 0.55] + [0.961] * 9
for ci, w in enumerate(col_widths):
    new_tbl.columns[ci].width = Inches(w)

# Header row: keep "KPI"/"unit", set the 9 WOBJ-config column headers.
for ci, (a, b) in enumerate(((a, b) for a in WOBJ_ALPHAS for b in WOBJ_BETAS), start=2):
    _set_cell_lines(new_tbl.cell(0, ci), [f"α={a:.1f}", f"β={b:.1f}"], size_pt=8, bold=True)

# Data rows: cols 0-1 (KPI/unit) are already correct via the clone; fill cols 2-10.
for ri, (label, unit, values) in enumerate(full_rows_wobj, start=1):
    for ci, val in enumerate(values, start=2):
        _set_cell(new_tbl.cell(ri, ci), val, size_pt=8, bold=False)

# Uniform 8pt font across the whole new table (cols 0-1 were cloned at 9pt).
for ri in range(N_ROWS):
    for ci in range(N_COLS2):
        cell = new_tbl.cell(ri, ci)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)

# Reposition the new slide immediately after the KPI-table slide (index 6 -> 7).
sldIdLst = prs.slides._sldIdLst
new_sldId = sldIdLst[-1]
sldIdLst.remove(new_sldId)
sldIdLst.insert(7, new_sldId)

# =============================================================================
# 4. Slide "1.2" -- configuration-count caption
# =============================================================================
s_sweeps = prs.slides[4]
caption_box = s_sweeps.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(12.13), Inches(1.35))
tf = caption_box.text_frame
tf.word_wrap = True
run = tf.paragraphs[0].add_run()
run.text = (
    "This period's dataset: 58 configurations run across 7 experiment families — "
    "6 main strategies (NoPriority + 5 coordinated DCTSP variants) plus the "
    "ZIG / BUSPRIO / WOBJ / MAINSIDE / LATENESS / TRAVELTIME sensitivity sweeps "
    "(see the following 2 slides for KPI tables on 6 and 9 of these configurations). "
    "The 3 sweeps above add 7 further configurations (4 detection-level + 3 "
    "occupancy-level, targeting NashGate) — configured, pending run. A demand-level "
    "sweep mode (6 runs: NoPriority + NashGate × {0.8, 1.0, 1.2}) is also available, "
    "off by default. Total configured experiment set this period: 65 "
    "(+ 6 demand-mode runs when enabled)."
)
run.font.size = Pt(11)
run.font.italic = True

# =============================================================================
# Save
# =============================================================================
prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
print(f"Total slides: {len(prs.slides)}")
