"""Experiment: does changing the graphicFrame's overall shape height/top
(not per-row heights) change the rendered table height in PowerPoint COM
export? Shrink the Preview table's frame to 5.6in tall, starting at 0.75in
(right under the title), and re-export to check.
"""

from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
prs = Presentation(PPTX_PATH)

s10 = prs.slides[10]
tbl_shape = next(sh for sh in s10.shapes if sh.has_table)

tbl_shape.top = Inches(0.75)
tbl_shape.height = Inches(5.6)

prs.save(PPTX_PATH)
print("Saved", PPTX_PATH)
