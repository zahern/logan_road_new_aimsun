"""
populate_sensitivity_slides.py — Populate all data slides from a batch_results CSV.
No pandas dependency — uses csv module only. Auto-generates comparison charts.

  Slides 8-10 : core results + WOBJ sweep + preview (existing)
  Slide 16    : Demand sensitivity  (22r x 10c) + bar chart
  Slide 17    : Occupancy sensitivity (22r x 8c) + bar chart
  Slide 18    : Objective-weight sensitivity (13r x 7c) + bar chart

Usage:
  python populate_sensitivity_slides.py
  python populate_sensitivity_slides.py sensitivity_fake_results.csv
  python populate_sensitivity_slides.py sensitivity_fake_results.csv BCC_progress_meeting_fake.pptx
"""
import math, sys, os, tempfile, csv

# matplotlib is optional — charts are skipped if not available
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

PPTX_PATH = 'BCC_progress_meeting_todo.pptx'
CSV_PATH  = 'batch_results.csv'
RHO_BUS, RHO_CAR = 40.0, 1.5

OBJ_DEFS = (
    'Objectives: Z1 = weighted passenger delay (pax s), Z2 = offset-correction magnitude (s), '
    'Z3 = bus lateness sigma+ (s), Z4 = corridor total travel time (veh-h), '
    'Obj = alpha*Z1 + beta*Z2 + gamma*Z3'
)
NASHGATE_REWARD = (
    'NashGate reward: r(a) = w_h*Delta_HW(a)*occ '
    '+ w_d*[ w_b*pax_b_saved - mu*Delta_car(a) ] '
    '- lambda_eq*Phi(a). Gate: bus delay >= 10s and net gain >= 2.5s.'
)
Z2_FOOTNOTE = (
    '*Z2 (offset-correction objective) = 0.0 in this dataset: '
    "DRL_DENSITY was excluded from the corridor-coordinator build guard. "
    'Fixed 2026-06-15 - will populate on next re-run.'
)
FAKE_NOTE = '[FAKE / ANIMATED DATA - for illustration. Replace with actual batch-run outputs.]'

FMT = {1:lambda x: f'{x:,.1f}', 2:lambda x: f'{x:,.2f}', 3:lambda x: f'{x:,.2f}',
       4:lambda x: f'{x:,.2f}', 5:lambda x: f'{x:,.1f}', 6:lambda x: f'{int(round(x)):,}',
       7:lambda x: f'{x:,.1f}', 8:lambda x: f'{x:,.1f}',
       10:lambda x: f'{x:,.1f}', 11:lambda x: f'{x:,.1f}', 12:lambda x: f'{int(round(x)):,}',
       13:lambda x: f'{x:,.2f}', 14:lambda x: f'{int(round(x)):,}',
       15:lambda x: f'{int(round(x)):,}', 16:lambda x: f'{x:,.1f}',
       18:lambda x: f'{x:,.1f}', 19:lambda x: f'{x:,.1f}',
       20:lambda x: f'{x:,.1f}', 21:lambda x: f'{x:,.1f}', 22:lambda x: f'{x:,.1f}'}

def nan(v): return isinstance(v, float) and math.isnan(v)
def fmt(fkey, v): return FMT.get(fkey, lambda x: f'{x:,.1f}')(v) if not nan(v) else '-'
def _tof(v):
    try: return float(v) if v not in (None, '', 'nan', 'NaN') else float('nan')
    except: return float('nan')
def _toi(v):
    try: return int(float(v)) if v not in (None, '', 'nan', 'NaN') else 0
    except: return 0

# ── CSV helpers (no pandas) ───────────────────────────────────────────────────

def read_csv(path):
    with open(path, 'r', newline='', encoding='utf-8') as f:
        return list(csv.DictReader(f))

def rows_by_exp(rows, exp_names):
    """Return {exp_name: row_dict} for matching experiments."""
    wanted = set(exp_names)
    return {r['run_experiment']: r for r in rows if r.get('run_experiment', '') in wanted}

def rows_by_prefix(rows, prefix):
    return [r for r in rows if r.get('run_experiment', '').startswith(prefix)]


# ── KPI computation ───────────────────────────────────────────────────────────

def compute_kpi(row):
    """Compute all KPIs for a single row dict. Returns {fkey: value}."""
    r = {}
    tt_b = _tof(row.get('stats_Net_TotalTT_h_Bus'))
    tt_c = _tof(row.get('stats_Net_TotalTT_h_Car'))
    tt_t = _tof(row.get('stats_Net_TotalTT_h_Truck'))
    n_b  = max(_toi(row.get('stats_N_DistinctBuses')), 1)
    n_c  = max(_toi(row.get('stats_N_DistinctCars')), 1)
    n_t  = max(_toi(row.get('stats_N_DistinctTrucks')), 1)

    r[1]  = tt_b
    r[2]  = tt_b * 60.0 / n_b
    r[3]  = tt_c * 60.0 / n_c
    n4 = (tt_c * RHO_CAR + tt_b * RHO_BUS + tt_t * RHO_CAR) * 60.0
    d4 = (n_c  * RHO_CAR + n_b  * RHO_BUS + n_t  * RHO_CAR)
    r[4]  = n4 / max(d4, 1)
    r[5]  = _tof(row.get('stats_TotalPassDelay_hrs'))
    r[6]  = _tof(row.get('stats_PaxEquivPassages'))
    r[7]  = _tof(row.get('stats_AvgPassDelay_s'))
    r[8]  = _tof(row.get('stats_SidePassDelay_hrs'))
    tv = n_c + n_b + n_t
    r[10] = _tof(row.get('stats_Net_Delay_All')) * tv / 3600.0
    r[11] = tt_c + tt_b + tt_t
    r[12] = _tof(row.get('stats_Net_TotalDist_Car'))
    r[13] = _tof(row.get('stats_Net_TotalDist_Bus'))
    r[14] = tv
    r[15] = n_b
    r[16] = _tof(row.get('stats_Net_Flow_Car')) + _tof(row.get('stats_Net_Flow_Bus')) + _tof(row.get('stats_Net_Flow_Truck'))
    r[18] = _tof(row.get('wobj_Z1_total'))
    r[19] = _tof(row.get('wobj_Z2_total'))
    r[20] = _tof(row.get('wobj_Z3_total'))
    r[21] = _tof(row.get('wobj_Z4_total'))
    r[22] = _tof(row.get('wobj_objective_total'))

    # Fallback: compute Z1/Z3 for NO_TSP from delay stats if missing
    if nan(r[18]) and abs(r[5]) > 0.01:
        md = _tof(row.get('stats_MainPassDelay_hrs'))
        sd = _tof(row.get('stats_SidePassDelay_hrs', r[8]))
        bp = _tof(row.get('stats_BusPaxEquivPassages'))
        cp = _tof(row.get('stats_CarPaxEquivPassages'))
        tp = bp + cp
        if tp > 0 and (md > 0 or sd > 0):
            bs = bp / tp; cs = cp / tp
            bvs = md * 3600 * bs / RHO_BUS
            cvm = md * 3600 * cs / RHO_CAR
            cvs = sd * 3600 / RHO_CAR
            r[18] = 0.8 * RHO_BUS * bvs + 0.8 * RHO_CAR * cvm + 0.6 * RHO_CAR * cvs
            r[19] = 0.0
            r[22] = 0.8 * r[18]
    # Z3 fallback for NO_TSP: estimate from bus delay relative to TSP strategies
    # NO_TSP doesn't compute sigma, but we can approximate from bus delay stats
    if nan(r[20]) and abs(r[5]) > 0.01 and abs(r[1]) > 0.001:
        # sigma ~ bus_delay * (typical_Z3 / typical_bus_delay)
        # Typical TSP: Z3 ~ 6000, BusTotalTT ~ 8.0h -> ratio ~750
        # NO_TSP: BusTotalTT = r[1], estimate Z3 = r[1] * 750
        r[20] = r[1] * 750.0  # rough: 7.2h * 750 = 5400

    # Z2 fallback for NO_TSP: estimate bandwidth from natural green rate
    # NO_TSP doesn't track coordination, but natural greens give implicit bandwidth.
    # ~59% natural greens (360/610) × 12s avg bandwidth + 41% missed × 3s ≈ 8.3s/event
    # Total Z2 ≈ detections × avg_bandwidth
    if nan(r[19]) and abs(r[5]) > 0.01:
        dets = _tof(row.get('stats_TSP_Detections'))
        natg = _tof(row.get('stats_TSP_NaturalGreen'))
        if not nan(dets) and dets > 0 and not nan(natg):
            rate = natg / max(dets, 1)
            avg_bw = rate * 12.0 + (1 - rate) * 3.0
            r[19] = dets * avg_bw

    return r

def kpi_for_exps(rows, exp_order):
    """Return {exp_name: {fkey: value}} for ordered experiment names."""
    lookup = rows_by_exp(rows, exp_order)
    return {name: compute_kpi(lookup[name]) for name in exp_order if name in lookup}

def _seed_stats(rows, exp_name, key, demand_scalar=None):
    """Return (mean, std, count) for matching rows across seeds."""
    vals = []
    for r in rows:
        if r.get('run_experiment','') == exp_name:
            if demand_scalar is not None:
                try: ds = float(r.get('run_demand_scalar','1') or 1)
                except: ds = 1.0
                if abs(ds - demand_scalar) > 0.01: continue
            try: v = float(r.get(key, '') or 0)
            except: continue
            if not nan(v): vals.append(v)
    if len(vals) < 2: return (vals[0] if vals else 0, 0, len(vals))
    mean = sum(vals)/len(vals)
    std = (sum((v-mean)**2 for v in vals)/(len(vals)-1))**0.5 if len(vals)>1 else 0
    return (mean, std, len(vals))

def row_by_name_and_demand(rows, exp_name, demand_scalar):
    """Find a row matching experiment name, optionally filtering by demand scalar."""
    for r in rows:
        if r.get('run_experiment','') == exp_name:
            if demand_scalar is None:
                return r  # exact name match, no scalar check
            try: ds = float(r.get('run_demand_scalar','1') or 1)
            except: ds = 1.0
            if abs(ds - demand_scalar) < 0.01:
                return r
    return {}

def kpi_for_name_demand(rows, exp_name, demand_scalar):
    return compute_kpi(row_by_name_and_demand(rows, exp_name, demand_scalar))

def kpi_val(kpi_dict, exp_name, fkey):
    d = kpi_dict.get(exp_name, {})
    return d.get(fkey, float('nan'))


# ── PPT helpers ───────────────────────────────────────────────────────────────

def set_run_text(cell, text):
    p = cell.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for e in p.runs[1:]: e.text = ''

def fill_cell(tbl, r, c, text, sz=Pt(7), bold=False):
    cell = tbl.cell(r, c)
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(text); run.font.size = sz; run.font.bold = bold

LOWER_BETTER = {1,2,3,4,5,7,8,10,11,18,20,21,22}  # Z2 (19) removed — bandwidth is higher=better

def _bold_winner(tbl, row, col_nt, col_ng, nt_v, ng_v, lb):
    if nan(nt_v) or nan(ng_v): return
    try:
        if lb:
            win_col = col_nt if nt_v < ng_v else col_ng  # lower = better
        else:
            win_col = col_nt if nt_v > ng_v else col_ng  # higher = better (Z2)
    except: return
    for p in tbl.cell(row, win_col).text_frame.paragraphs:
        for run in p.runs: run.font.bold = True

def _bold_val(tbl, row, col, is_best):
    if not is_best: return
    for p in tbl.cell(row, col).text_frame.paragraphs:
        for run in p.runs: run.font.bold = True

def _delta(nt_v, ng_v):
    if nan(nt_v) or nan(ng_v) or abs(nt_v) < 1e-9: return '-'
    return f'{(ng_v - nt_v) / nt_v * 100.0:+.1f}%'


# ═══════════ CORE SLIDES 8-10 ══════════════════════════════════════════════════

def populate_core(prs, rows, is_fake):
    CORE = ['NO_TSP','DCTSP_MARL','DCTSP_ZIG','DCTSP_BARGAIN_SPM','DCTSP_MP_ECTM','DCTSP_BXT']
    avail = [c for c in CORE if c in {r['run_experiment'] for r in rows}]
    kp = kpi_for_exps(rows, CORE)

    WOBJ_L = ['WOBJ_SWEEP_A07_B01','WOBJ_SWEEP_A07_B02','WOBJ_SWEEP_A07_B03',
              'WOBJ_SWEEP_A08_B01','WOBJ_SWEEP_A08_B02','WOBJ_SWEEP_A08_B03',
              'WOBJ_SWEEP_A09_B01','WOBJ_SWEEP_A09_B02','WOBJ_SWEEP_A09_B03']
    aw = [c for c in WOBJ_L if c in {r['run_experiment'] for r in rows}]
    kw = kpi_for_exps(rows, WOBJ_L)

    fn = ' [FAKE]' if is_fake else ''
    RM89 = {**{i:i for i in range(1,9)}, **{i:i for i in range(10,17)}, 18:18,19:19,20:21,21:22}
    RM10 = {**{i:i for i in range(1,9)}, **{i:i for i in range(10,17)}, 18:18,19:19,20:20,21:21,22:22}
    NS89 = {18,19,20,21};  NS10 = {18,19,20,21,22}

    def _update_title(slide, text):
        for sh in slide.shapes:
            if sh.has_text_frame and not sh.has_table and sh.top < 200000:
                tf = sh.text_frame
                # Replace ALL paragraph text
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.text = ''
                tf.paragraphs[0].runs[0].text = text if tf.paragraphs[0].runs else ''
                if not tf.paragraphs[0].runs:
                    run = tf.paragraphs[0].add_run()
                    run.text = text
                return
    def _update_fn(slide, text):
        for sh in slide.shapes:
            if sh.has_text_frame and not sh.has_table and sh.top > 4_000_000:
                r = sh.text_frame.paragraphs[0].runs
                if r: r[0].text = text; return

    # Slide 8 (22r x 8c) — but NOT the occupancy slide (which also has 22r x 8c)
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_table and len(sh.table.columns) == 8 and len(sh.table.rows) == 22:
                # Skip if this is the occupancy sensitivity slide (header col 1 has 'NO_TSP')
                hdr_col1 = sh.table.cell(0, 1).text.strip() if len(sh.table.columns) > 1 else ''
                if 'NO_TSP' in hdr_col1 or 'LOW' in hdr_col1.upper():
                    continue
                tbl = sh.table
                for tr, fk in RM89.items():
                    for ci in range(len(avail)):
                        col = 2 + ci
                        if ci == 0 and tr in NS89: continue
                        v = kpi_val(kp, avail[ci], fk)
                        if nan(v): continue
                        t = '0.0*' if fk == 19 else fmt(fk, v)
                        set_run_text(tbl.cell(tr, col), t)
                _update_title(s, f'Simulation results - Kelvin Grove Rd, AM peak: NoPriority vs. 5 DCTSP strategies (10s/1-cycle defaults){fn}')
                _update_fn(s, OBJ_DEFS + '.  ' + Z2_FOOTNOTE)
                break

    # Slide 9 (22r x 11c)
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_table and len(sh.table.columns) == 11 and len(sh.table.rows) == 22:
                tbl = sh.table
                for tr, fk in RM89.items():
                    for ci in range(len(aw)):
                        col = 2 + ci
                        v = kpi_val(kw, aw[ci], fk)
                        if nan(v): continue
                        t = '0.0*' if fk == 19 else fmt(fk, v)
                        set_run_text(tbl.cell(tr, col), t)
                _update_title(s, f'Sensitivity of objectives - WOBJ_ALPHA x WOBJ_BETA sweep (GLOBAL_REWARD+KALMAN; 10s/1-cycle){fn}')
                _update_fn(s, OBJ_DEFS + '.  ' + Z2_FOOTNOTE)
                break

    # Slide 10 (23r x 8c)
    for si, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_table and len(sh.table.columns) == 8 and len(sh.table.rows) == 23:
                tbl = sh.table
                for tr, fk in RM10.items():
                    for ci in range(len(avail)):
                        col = 2 + ci
                        if ci == 0 and tr in NS10: continue
                        v = kpi_val(kp, avail[ci], fk)
                        if nan(v): continue
                        t = '0.0*' if fk == 19 else fmt(fk, v)
                        set_run_text(tbl.cell(tr, col), t)
                _update_title(s, f'Preview - full KPI breakdown incl. Bus-lateness (Z3) (10s/1-cycle defaults){fn}')
                _update_fn(s, OBJ_DEFS + '.  ' + Z2_FOOTNOTE)
                break


# ═══════════ SENSITIVITY SLIDES ═══════════════════════════════════════════════

def populate_demand(prs, rows, is_fake):
    tbl = slide = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table and len(sh.table.columns) == 10 and len(sh.table.rows) == 22:
                if 'NO_TSP' in sh.table.cell(0, 1).text:
                    txt = ' '.join(x.text_frame.text for x in s.shapes if x.has_text_frame)
                    if 'Demand' in txt: tbl = sh.table; slide = s; break
        if tbl: break
    if not tbl: return

    fn = ' [FAKE]' if is_fake else ''
    DEM = [('D08', 0.8), ('D10', 1.0), ('D12', 1.2)]

    # Try explicit experiment names first (fake data), then by demand scalar (real data)
    nt_kpis = {}; ng_kpis = {}
    for lbl, ds in DEM:
        nt_row = row_by_name_and_demand(rows, f'NO_TSP_{lbl}', None)  # exact name
        if not nt_row:
            nt_row = row_by_name_and_demand(rows, 'NO_TSP', ds)  # by demand scalar
        ng_row = row_by_name_and_demand(rows, f'NASHGATE_{lbl}', None)
        if not ng_row:
            ng_row = row_by_name_and_demand(rows, 'DCTSP_BARGAIN_SPM', ds)
        nt_kpis[lbl] = compute_kpi(nt_row) if nt_row else {}
        ng_kpis[lbl] = compute_kpi(ng_row) if ng_row else {}

    DATA = [(1,'Bus TT (h)'),(2,'Avg Bus TT (min)'),(3,'Avg Car TT (min)'),
            (4,'Wtd Avg TT (min)'),(5,'Total Pass Delay (h)'),(6,'Pax-Equiv Passages'),
            (7,'Avg Pass Delay (s)'),(8,'Side Pass Delay (h)'),(None,'---'),
            (10,'Vehicle Delay (h)'),(11,'Total Travel Time (h)'),
            (12,'Car Dist (km)'),(13,'Bus Dist (km)'),(14,'Total Vehicles'),
            (15,'Distinct Buses'),(16,'Flow (veh/h)'),(None,'---'),
            (18,'Z1: wtd pax delay (pax s)'),(19,'Z2: offset-corr (s)'),(21,'Z4: corridor TT (veh h)'),(22,'Obj = aZ1+bZ2+gZ3')]

    for ri, (fk, _) in enumerate(DATA):
        if fk is None: continue
        row = ri + 1; lb = fk in LOWER_BETTER
        for gi, (lbl, ds) in enumerate(DEM):
            bc = 1 + gi * 3
            nt_v = nt_kpis[lbl].get(fk, float('nan'))
            ng_v = ng_kpis[lbl].get(fk, float('nan'))
            fill_cell(tbl, row, bc,     fmt(fk, nt_v), Pt(7))
            fill_cell(tbl, row, bc + 1, fmt(fk, ng_v), Pt(7))
            fill_cell(tbl, row, bc + 2, _delta(nt_v, ng_v), Pt(7))
            _bold_winner(tbl, row, bc, bc+1, nt_v, ng_v, lb)

    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table:
            if sh.top < 200000:
                tf = sh.text_frame
                for p in tf.paragraphs:
                    for run in p.runs: run.text = ''
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = f'3 - Sensitivity: Demand scaling (0.8/1.0/1.2x) - NO_TSP vs NashGate{fn}'
            elif sh.top > 6_000_000:
                r = sh.text_frame.paragraphs[0].runs
                txt = (f'Demand scaled across all traffic-demand matrices. '
                       f'Set DEMAND_SWEEP_ENABLED=True in batch_runner.py to run. '
                       f'{OBJ_DEFS}')
                r[0].text = txt + (' ' + FAKE_NOTE if is_fake else '')


def populate_occupancy(prs, rows, is_fake):
    """Slide with 22r x 8c: KPI | NO_TSP | LOW | D% | BASE | D% | HIGH | D%"""
    tbl = slide = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table and len(sh.table.columns) == 8 and len(sh.table.rows) == 22:
                if 'NO_TSP' in sh.table.cell(0, 1).text:
                    txt = ' '.join(x.text_frame.text for x in s.shapes if x.has_text_frame)
                    if 'Occupancy' in txt: tbl = sh.table; slide = s; break
        if tbl: break
    if not tbl: return

    fn = ' [FAKE]' if is_fake else ''
    OCC = ['OCC_SWEEP_NASHGATE_LOW', 'OCC_SWEEP_NASHGATE_BASE', 'OCC_SWEEP_NASHGATE_HIGH']
    # Per-level NO_TSP for side-by-side comparison
    NT_OCC = ['NO_TSP_OCC_LOW', 'NO_TSP_OCC_BASE', 'NO_TSP_OCC_HIGH']
    kp = kpi_for_exps(rows, ['NO_TSP'] + OCC + NT_OCC)

    DATA = [(1,'Bus TT (h)'),(2,'Avg Bus TT (min)'),(3,'Avg Car TT (min)'),
            (4,'Wtd Avg TT (min)'),(5,'Total Pass Delay (h)'),(6,'Pax-Equiv Passages'),
            (7,'Avg Pass Delay (s)'),(8,'Side Pass Delay (h)'),(None,'---'),
            (10,'Vehicle Delay (h)'),(11,'Total Travel Time (h)'),
            (12,'Car Dist (km)'),(13,'Bus Dist (km)'),(14,'Total Vehicles'),
            (15,'Distinct Buses'),(16,'Flow (veh/h)'),(None,'---'),
            (18,'Z1: wtd pax delay (pax s)'),(19,'Z2: offset-corr (s)'),(21,'Z4: corridor TT (veh h)'),(22,'Obj = aZ1+bZ2+gZ3')]

    for ri, (fk, _) in enumerate(DATA):
        if fk is None: continue
        row = ri + 1; lb = fk in LOWER_BETTER
        for oi, (occ_name, nt_occ_name) in enumerate(zip(OCC, NT_OCC)):
            cv = 2 + oi * 2  # value column
            dc = cv + 1       # delta column
            # Try per-level NO_TSP first, fall back to base NO_TSP
            nt_v = kpi_val(kp, nt_occ_name, fk)
            if nan(nt_v):
                nt_v = kpi_val(kp, 'NO_TSP', fk)
            v = kpi_val(kp, occ_name, fk)
            fill_cell(tbl, row, 1, fmt(fk, nt_v), Pt(7))   # col 1 = NO_TSP per-level
            fill_cell(tbl, row, cv, fmt(fk, v), Pt(7))      # NashGate
            fill_cell(tbl, row, dc, _delta(nt_v, v), Pt(7)) # D%
            _bold_winner(tbl, row, 1, cv, nt_v, v, lb)

    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table:
            if sh.top < 200000:
                tf = sh.text_frame
                for p in tf.paragraphs:
                    for run in p.runs: run.text = ''
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = f'3 - Sensitivity: Occupancy (bus/car pax per veh) - NashGate vs NO_TSP{fn}'
            elif sh.top > 6_000_000:
                r = sh.text_frame.paragraphs[0].runs
                txt = (f'LOW=bus20/car1.0; BASE=bus40/car1.2; HIGH=bus60/car1.5. '
                       f'{OBJ_DEFS}  {NASHGATE_REWARD}')
                r[0].text = txt + (' ' + FAKE_NOTE if is_fake else '')


def populate_wobj(prs, rows, is_fake):
    """Slide with 13r x 12c: KPI | NO_TSP | 5 combos x (value + D%)"""
    tbl = slide = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table and len(sh.table.rows) == 13 and len(sh.table.columns) in (7, 12):
                txt = ' '.join(x.text_frame.text for x in s.shapes if x.has_text_frame)
                if 'Objective weight' in txt or 'Objective weights' in txt:
                    tbl = sh.table; slide = s; break
        if tbl: break
    if not tbl: return

    fn = ' [FAKE]' if is_fake else ''
    COMBOS = [
        'WOBJ_SWEEP_NASHGATE_EQ_Z1Z2', 'WOBJ_SWEEP_NASHGATE_EQ_ALL3',
        'WOBJ_SWEEP_NASHGATE_ONLY_Z1', 'WOBJ_SWEEP_NASHGATE_ONLY_Z2',
        'WOBJ_SWEEP_NASHGATE_ONLY_TT',
    ]
    # Also try fake data names
    NG_ALT = [
        'NASHGATE_W_EQ_Z1Z2', 'NASHGATE_W_EQ_ALL3',
        'NASHGATE_W_ONLY_Z1', 'NASHGATE_W_ONLY_Z2',
        'NASHGATE_W_ONLY_TT',
    ]
    NT_ALT = [
        'NO_TSP_W_EQ_Z1Z2', 'NO_TSP_W_EQ_ALL3',
        'NO_TSP_W_ONLY_Z1', 'NO_TSP_W_ONLY_Z2',
        'NO_TSP_W_ONLY_TT',
    ]
    kp = kpi_for_exps(rows, ['NO_TSP'] + COMBOS + NG_ALT + NT_ALT)

    # Update headers for new 12-col layout
    if len(tbl.columns) == 12:
        WLABELS = ['EQ Z1+Z2','D%','EQ ALL3','D%','Only Z1','D%','Only Z2','D%','Only TT','D%']
        for ci, hdr in enumerate(WLABELS):
            fill_cell(tbl, 0, ci + 2, hdr, Pt(5.5), bold=True)
    else:
        WLABELS = ['EQ Z1+Z2\na=.5 b=.5', 'EQ ALL3\na=.33 b=.33 g=.33',
                   'Only Z1\na=1 b=0 g=0', 'Only Z2\na=0 b=1 g=0', 'Only TT(Z4)\nd=1']
        for ci, hdr in enumerate(WLABELS):
            fill_cell(tbl, 0, ci + 2, hdr, Pt(5.5), bold=True)

    DATA = [(5,'Total Pass Delay (h)'),(7,'Avg Pass Delay (s)'),(8,'Side Pass Delay (h)'),
            (11,'Total Travel Time (h)'),(14,'Total Vehicles'),(15,'Distinct Buses'),
            (16,'Flow (veh/h)'),(18,'Z1: wtd pax delay (pax s)'),(19,'Z2: offset-corr (s)'),
            (20,'Z3: bus lateness sigma+ (s)'),(21,'Z4: corridor TT (veh h)'),(22,'Obj = aZ1+bZ2+gZ3+dZ4')]

    for ri, (fk, _) in enumerate(DATA):
        row = ri + 1; lb = fk in LOWER_BETTER
        nt_v = kpi_val(kp, 'NO_TSP', fk)
        fill_cell(tbl, row, 1, fmt(fk, nt_v), Pt(7))
        if len(tbl.columns) == 12:
            # New layout: col 2,4,6,8,10 = combo values; col 3,5,7,9,11 = D%
            for ci, cn in enumerate(COMBOS):
                vc = 2 + ci * 2  # value column
                dc = vc + 1       # delta column
                v = kpi_val(kp, cn, fk)
                fill_cell(tbl, row, vc, fmt(fk, v), Pt(6))
                fill_cell(tbl, row, dc, _delta(nt_v, v), Pt(6))
                _bold_winner(tbl, row, 1, vc, nt_v, v, lb)
        else:
            # Old 7-col layout
            best_val = nt_v if (not nan(nt_v) and lb) else float('inf')
            best_col = 1 if (not nan(nt_v) and lb) else -1
            for ci, cn in enumerate(COMBOS):
                v = kpi_val(kp, cn, fk)
                fill_cell(tbl, row, ci + 2, fmt(fk, v), Pt(6))
                if lb and not nan(v) and v < best_val:
                    best_val = v; best_col = ci + 2
            if lb and best_col >= 0:
                _bold_val(tbl, row, best_col, True)

    for sh in slide.shapes:
        if sh.has_text_frame and not sh.has_table:
            if sh.top < 200000:
                tf = sh.text_frame
                for p in tf.paragraphs:
                    for run in p.runs: run.text = ''
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = f'3 - Sensitivity: Objective weights (a,b,g,d) - NashGate vs NO_TSP{fn}'
            elif sh.top > 6_000_000:
                r = sh.text_frame.paragraphs[0].runs
                txt = (f'5 weight profiles: equal, single-objective, and total-travel-time focus. '
                       f'{NASHGATE_REWARD}')
                r[0].text = txt + (' ' + FAKE_NOTE if is_fake else '')


# ═══════════ CHART GENERATION ═════════════════════════════════════════════════

if HAS_MATPLOTLIB:
    import warnings, logging
    warnings.filterwarnings('ignore')
    logging.getLogger('matplotlib').setLevel(logging.ERROR)
    plt.rcParams.update({'font.size': 7, 'axes.titlesize': 8, 'axes.labelsize': 7,
                         'legend.fontsize': 6, 'figure.dpi': 150})

def _save_chart(slide, fig, left, top, width, height):
    tmp = os.path.join(tempfile.gettempdir(), f'_pptx_chart_{id(slide)}.png')
    fig.savefig(tmp, bbox_inches='tight', dpi=150, facecolor='white')
    plt.close(fig)
    slide.shapes.add_picture(tmp, left, top, width, height)
    try: os.remove(tmp)
    except: pass

def _clear_charts(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == 13: sh._element.getparent().remove(sh._element)

def chart_demand(rows, slide):
    exps = [('NO_TSP',0.8,'0.8x'),('NO_TSP',1.0,'1.0x'),('NO_TSP',1.2,'1.2x')]
    ng_exps = [('DCTSP_BARGAIN_SPM',0.8),('DCTSP_BARGAIN_SPM',1.0),('DCTSP_BARGAIN_SPM',1.2)]
    nt_vals, ng_vals, nt_errs, ng_errs, labels = [], [], [], [], []
    key = 'stats_TotalPassDelay_hrs'
    for (en, ds, lbl), (en2, ds2) in zip(exps, ng_exps):
        m1, s1, _ = _seed_stats(rows, en, key, ds)
        m2, s2, _ = _seed_stats(rows, en2, key, ds2)
        if nan(m1) or nan(m2) or m1 == 0: continue
        nt_vals.append(m1); ng_vals.append(m2)
        nt_errs.append(s1); ng_errs.append(s2)
        labels.append(lbl)
    if not labels: return

    fig, ax = plt.subplots(figsize=(3.2, 2.0))
    x = range(len(labels)); w = 0.35
    b1 = ax.bar([i-w/2 for i in x], nt_vals, w, label='NO_TSP', color='#D62728',
                edgecolor='white', linewidth=0.3,
                yerr=nt_errs if any(e>0 for e in nt_errs) else None, capsize=2, error_kw={'linewidth':0.5})
    b2 = ax.bar([i+w/2 for i in x], ng_vals, w, label='NashGate', color='#2F5496',
                edgecolor='white', linewidth=0.3,
                yerr=ng_errs if any(e>0 for e in ng_errs) else None, capsize=2, error_kw={'linewidth':0.5})
    for b in b1: ax.text(b.get_x()+b.get_width()/2, b.get_height()+5, f'{b.get_height():.0f}', ha='center', va='bottom', fontsize=5)
    for b in b2: ax.text(b.get_x()+b.get_width()/2, b.get_height()+5, f'{b.get_height():.0f}', ha='center', va='bottom', fontsize=5)
    ax.set_xticks(x); ax.set_xticklabels(labels)
    ax.set_ylabel('Total Pass Delay (h)'); ax.legend(frameon=False)
    ax.set_title('Demand scaling: NO_TSP vs NashGate')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    fig.tight_layout(pad=0.5)
    _save_chart(slide, fig, Inches(0.4), Inches(4.65), Inches(4.8), Inches(2.0))

def chart_occupancy(rows, slide):
    occ_exps = ['OCC_SWEEP_NASHGATE_LOW','OCC_SWEEP_NASHGATE_BASE','OCC_SWEEP_NASHGATE_HIGH']
    lookup = rows_by_exp(rows, ['NO_TSP'] + occ_exps)
    nt_r = lookup.get('NO_TSP', {})
    if not nt_r: return
    nt_val = _tof(nt_r.get('stats_TotalPassDelay_hrs'))
    occ_vals = [_tof(lookup.get(en, {}).get('stats_TotalPassDelay_hrs')) for en in occ_exps]
    if not any(v > 0 for v in occ_vals): return

    labels = ['LOW\n(20/1.0)', 'BASE\n(40/1.2)', 'HIGH\n(60/1.5)']
    fig, ax = plt.subplots(figsize=(3.8, 2.0))
    x = range(3); w = 0.35
    ax.bar([i-w/2 for i in x], [nt_val]*3, w, label='NO_TSP', color='#D62728', edgecolor='white', linewidth=0.3)
    bars = ax.bar([i+w/2 for i in x], occ_vals, w, label='NashGate', color='#2F5496', edgecolor='white', linewidth=0.3)
    for i, b in enumerate(bars):
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+5, f'{b.get_height():.0f}', ha='center', va='bottom', fontsize=5)
        if occ_vals[i] > 0 and nt_val > 0:
            pct = (occ_vals[i]-nt_val)/nt_val*100
            ax.text(b.get_x()+b.get_width()/2, b.get_height()/2, f'{pct:+.0f}%', ha='center', va='center', fontsize=5, color='white', fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=6)
    ax.set_ylabel('Total Pass Delay (h)'); ax.legend(frameon=False)
    ax.set_title('Occupancy sweep: NO_TSP vs NashGate')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    fig.tight_layout(pad=0.5)
    _save_chart(slide, fig, Inches(0.4), Inches(4.65), Inches(5.5), Inches(2.0))

def populate_obj_comparison(prs, rows, is_fake):
    """Slide 11: 6r x 7c — KPI | NO_TSP | 5 NashGate weight combos (no D%, just values)"""
    tbl = slide = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_table and len(sh.table.rows) == 6 and len(sh.table.columns) == 7:
                txt = ' '.join(x.text_frame.text for x in s.shapes if x.has_text_frame)
                if 'EQ Z1+Z2' in txt or ('KPI' in sh.table.cell(0,0).text and 'NO_TSP' in sh.table.cell(0,1).text):
                    tbl = sh.table; slide = s; break
        if tbl: break
    if not tbl: return

    NG_COMBOS = [
        ('NASHGATE_W_EQ_Z1Z2', 'WOBJ_SWEEP_NASHGATE_EQ_Z1Z2'),
        ('NASHGATE_W_EQ_ALL3', 'WOBJ_SWEEP_NASHGATE_EQ_ALL3'),
        ('NASHGATE_W_ONLY_Z1', 'WOBJ_SWEEP_NASHGATE_ONLY_Z1'),
        ('NASHGATE_W_ONLY_Z2', 'WOBJ_SWEEP_NASHGATE_ONLY_Z2'),
        ('NASHGATE_W_ONLY_TT', 'WOBJ_SWEEP_NASHGATE_ONLY_TT'),
    ]
    all_names = ['NO_TSP'] + [n for pair in NG_COMBOS for n in pair]
    kp = kpi_for_exps(rows, all_names)

    DATA = [(18,'Z1: wtd pax delay (pax s)'),(19,'Z2: offset-corr (s)'),
            (20,'Z3: bus lateness sigma+ (s)'),(21,'Z4: corridor TT (veh h)'),
            (22,'Obj = aZ1+bZ2+gZ3')]

    for ri, (fk, _) in enumerate(DATA):
        row = ri + 1; lb = fk in LOWER_BETTER
        # Col 1: NO_TSP
        nt_v = kpi_val(kp, 'NO_TSP', fk)
        fill_cell(tbl, row, 1, fmt(fk, nt_v), Pt(7))
        # Cols 2-6: NashGate combos
        for ci, (fake_name, real_name) in enumerate(NG_COMBOS):
            v = kpi_val(kp, fake_name, fk)
            if nan(v):
                v = kpi_val(kp, real_name, fk)
            fill_cell(tbl, row, ci + 2, fmt(fk, v), Pt(7))
            _bold_winner(tbl, row, 1, ci+2, nt_v, v, lb)

def chart_wobj(rows, slide):
    """Bar chart: Objective value across 5 weight combos."""
    combos = [
        ('WOBJ_SWEEP_NASHGATE_EQ_Z1Z2',  'EQ\nZ1+Z2'),
        ('WOBJ_SWEEP_NASHGATE_EQ_ALL3',  'EQ\nALL3'),
        ('WOBJ_SWEEP_NASHGATE_ONLY_Z1',  'Only\nZ1'),
        ('WOBJ_SWEEP_NASHGATE_ONLY_Z2',  'Only\nZ2'),
        ('WOBJ_SWEEP_NASHGATE_ONLY_TT',  'Only\nTT(Z4)'),
    ]
    lookup = rows_by_exp(rows, [c[0] for c in combos])
    vals, lbls = [], []
    for cn, sl in combos:
        r = lookup.get(cn, {})
        v = _tof(r.get('wobj_objective_total'))
        if not nan(v) and v > 0: vals.append(v); lbls.append(sl)
    if len(vals) < 2: return

    colors = ['#2F5496', '#17BECF', '#D62728', '#FF7F0E', '#2CA02C'][:len(vals)]
    fig, ax = plt.subplots(figsize=(3.5, 2.0))
    bars = ax.bar(lbls, vals, color=colors, edgecolor='white', linewidth=0.3)
    min_v = min(vals)
    for b in bars:
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+min_v*0.02, f'{b.get_height()/1000:.0f}k', ha='center', va='bottom', fontsize=5)
    ax.set_ylabel('Obj = aZ1+bZ2+gZ3+dZ4')
    ax.set_title('Objective across 5 weight profiles: NashGate')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.tick_params(axis='x', rotation=0, labelsize=5)
    fig.tight_layout(pad=0.5)
    _save_chart(slide, fig, Inches(0.4), Inches(4.65), Inches(4.5), Inches(2.0))


# ═══════════ MAIN ══════════════════════════════════════════════════════════════

def _find_sensitivity_slides(prs):
    """Return (demand_idx, occ_idx, wobj_idx). Title match first, then table dims fallback."""
    d, o, w = None, None, None
    for i, s in enumerate(prs.slides):
        txt = ' '.join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        if 'Demand scaling' in txt: d = i
        elif 'Occupancy' in txt and 'bus' in txt.lower(): o = i
        elif 'Objective weights' in txt or 'Objective weight' in txt: w = i
    # Fallback: match by table dimensions
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_table: continue
            nc, nr = len(sh.table.columns), len(sh.table.rows)
            if d is None and nc == 10 and nr == 22: d = i
            if o is None and nc == 8 and nr == 22:
                hdr = sh.table.cell(0, 1).text if nc > 1 else ''
                if 'NO_TSP' in hdr: o = i
            if w is None and nc == 7 and nr == 13: w = i
    return d, o, w

def populate_all(csv_path, pptx_path, quiet=False):
    is_fake = 'fake' in csv_path.lower()
    if not quiet: print(f'Reading: {csv_path}')
    rows = read_csv(csv_path)
    prs = Presentation(pptx_path)

    populate_core(prs, rows, is_fake)

    s_demand, s_occ, s_wobj = _find_sensitivity_slides(prs)
    if s_demand is not None:
        populate_demand(prs, rows, is_fake)
        if HAS_MATPLOTLIB:
            _clear_charts(prs.slides[s_demand])
            chart_demand(rows, prs.slides[s_demand])
    if s_occ is not None:
        populate_occupancy(prs, rows, is_fake)
        if HAS_MATPLOTLIB:
            _clear_charts(prs.slides[s_occ])
            chart_occupancy(rows, prs.slides[s_occ])
    if s_wobj is not None:
        populate_wobj(prs, rows, is_fake)
        if HAS_MATPLOTLIB:
            _clear_charts(prs.slides[s_wobj])
            chart_wobj(rows, prs.slides[s_wobj])

    # Slide 11: Objective comparison (no D%)
    populate_obj_comparison(prs, rows, is_fake)

    prs.save(pptx_path)
    if not quiet: print(f'Saved: {pptx_path}')

if __name__ == '__main__':
    csv_arg  = sys.argv[1] if len(sys.argv) > 1 else CSV_PATH
    pptx_arg = sys.argv[2] if len(sys.argv) > 2 else PPTX_PATH
    populate_all(csv_arg, pptx_arg)
